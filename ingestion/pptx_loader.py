from pptx import Presentation
from loguru import logger


def load_pptx(file_path: str) -> str:
    presentation = Presentation(file_path)
    text = ""

    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    if paragraph.text.strip():
                        text += paragraph.text + "\n"

    logger.info(f"PPTX loaded: {file_path}")
    return text