import sys
import os
from unittest import result
os.environ["TOKENIZERS_PARALLELISM"] = "false"
os.environ["TRANSFORMERS_VERBOSITY"] = "error"
os.environ["HF_HUB_DISABLE_PROGRESS_BARS"] = "1"
os.environ["HF_HUB_DISABLE_TELEMETRY"] = "1"
os.environ["PYTHONIOENCODING"] = "utf-8"
sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from dotenv import load_dotenv
load_dotenv()

import concurrent.futures
import io
import json
import re
import threading
import time

import streamlit as st
from loguru import logger

import sys as _sys
logger.remove()
logger.add(_sys.stderr, level="WARNING")

import os as _os
if _os.path.isdir("logs") or not _os.path.exists("logs"):
    try:
        _os.makedirs("logs", exist_ok=True)
        logger.add("logs/documind.log", level="DEBUG", rotation="10 MB", retention="7 days", enqueue=True)
    except Exception:
        pass

st.set_page_config(page_title="DocuMind AI", page_icon="DM", layout="wide", initial_sidebar_state="expanded")


@st.cache_resource(show_spinner=False)
def load_core():
    from ingestion.document_loader import load_document
    from chunking.text_chunker import chunk_text
    from vector_store.faiss_store import create_vector_store
    from vector_store.bm25_store import create_bm25_index
    from graph.workflow import run_workflow
    from memory.session_memory import SessionMemory
    from memory.file_memory_manager import FileMemoryManager
    from utils.validator import validate_file, validate_question

    # Background hardcoding guard — daemon thread, never blocks the UI
    def _run_guard():
        try:
            from utils.validator import scan_source_files
            root       = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
            violations = scan_source_files(root)
            if violations:
                logger.warning(f"HARDCODING GUARD: {len(violations)} issue(s) found:")
                for fpath, lineno, line_text, label in violations:
                    logger.warning(f"  {fpath}:{lineno} [{label}] → {line_text[:100]}")
            else:
                logger.debug("HARDCODING GUARD: clean.")
        except Exception as e:
            logger.debug(f"HARDCODING GUARD skipped: {e}")

    threading.Thread(target=_run_guard, daemon=True).start()

    # Pre-warm embedding model while user is on the upload screen
    from embeddings.embedding_model import get_embedding_model as _gw
    threading.Thread(target=_gw, daemon=True).start()

    return dict(
        load_document=load_document,
        chunk_text=chunk_text,
        create_vector_store=create_vector_store,
        create_bm25_index=create_bm25_index,
        run_workflow=run_workflow,
        SessionMemory=SessionMemory,
        FileMemoryManager=FileMemoryManager,
        validate_file=validate_file,
        validate_question=validate_question,
    )


@st.cache_resource(show_spinner=False)
def _llm_groq():
    key = os.getenv("GROQ_API_KEY", "")
    if key and key not in {"your_groq_key_here", ""}:
        try:
            from langchain_groq import ChatGroq
            return ChatGroq(
                model="llama-3.3-70b-versatile",
                groq_api_key=key,
                temperature=0.3,
                max_tokens=2048,
            )
        except Exception:
            pass
    return None


@st.cache_resource(show_spinner=False)
def _llm_gemini():
    from langchain_google_genai import ChatGoogleGenerativeAI
    from utils.config import get_google_api_key
    return ChatGoogleGenerativeAI(
        model="models/gemini-2.5-flash",
        google_api_key=get_google_api_key(),
        temperature=0.3,
    )


def get_llm(temp=0.3):
    llm = _llm_groq() or _llm_gemini()
    if temp != 0.3:
        try:
            return llm.with_config(configurable={"temperature": temp})
        except Exception:
            pass
    return llm


def lazy(mod, fn):
    import importlib
    return getattr(importlib.import_module(mod), fn)


def uid():
    return f"{int(time.time() * 1000)}_{len(st.session_state.get('chat', []))}"


def init(deps):
    simple = {
        "chat": [],
        "file_path": "",
        "doc_ready": False,
        "doc_name": "",
        "doc_text": "",
        "prompts": [],
        "answer_mode": "detailed",
        "multi_docs": [],
        "last_q": "",
        "last_a": "",
        "pending": {},
        "quiz_store": {},
        "is_processing": False,
        "stop_requested": False,
        "_processing": False,
        "_stop_flag": False,
        "startup_checked": False,
        "startup_results": {},
    }
    for k, v in simple.items():
        if k not in st.session_state:
            st.session_state[k] = v
    if "memory" not in st.session_state:
        st.session_state["memory"] = deps["SessionMemory"]()
    if "fmm" not in st.session_state:
        st.session_state["fmm"] = deps["FileMemoryManager"]()


def detect_intent(q):
    import re as _re
    ql = q.lower()
    toks = set(_re.findall(r"[a-z0-9]+", ql))

    # Complaint/follow-up: do not trigger another document generation.
    _complaints = [
        "you provided only", "you gave only", "you only gave",
        "you only provided", "you provided", "you gave",
        "you said", "you told me", "you mentioned",
        "that's wrong", "that is wrong", "not right", "incorrect",
        "you made a mistake", "not what i asked", "that's not what",
        "you only showed", "you only sent",
    ]
    if any(p in ql for p in _complaints):
        return {k: False for k in
                ["chart", "quiz", "excel", "pptx", "export_doc",
                 "resources", "regen_doc"]}

    action_toks = {
        "generate", "create", "make", "prepare", "build",
        "provide", "give", "get", "export", "produce",
        "download", "show",
    }
    doc_toks = {
        "document", "doc", "docx", "pdf", "file", "report", "sheet",
    }
    wants_doc = bool(toks & action_toks) and bool(toks & doc_toks)
    has_n_records = bool(_re.search(
        r"\b\d+\s+(?:records?|rows?|entries|entry|people|persons?|respondents?|items?)\b",
        ql,
    ))

    return {
        "chart": any(w in ql for w in [
            "chart", "graph", "plot", "bar chart", "pie chart", "line graph",
            "visualize", "histogram", "draw a", "show graph", "show chart",
            "bar graph", "scatter", "visual",
        ]),
        "quiz": any(w in ql for w in [
            "quiz", "test me", "mcq", "multiple choice", "exam", "question me",
            "conduct a quiz", "give me a quiz",
        ]),
        "excel": (
            "excel" in toks
            or "xlsx" in toks
            or "generate excel" in ql
            or "export excel" in ql
            or "spreadsheet file" in ql
        ),
        "pptx": any(w in ql for w in [
            "powerpoint", "pptx", "presentation", "make slides", "create slides",
        ]),
        "export_doc": any(w in ql for w in [
            "export as", "save as", "generate pdf", "create pdf", "export pdf",
            "generate docx", "create docx", "export doc", "generate csv",
            "download as",
        ]),
        "resources": any(w in ql for w in [
            "resources", "youtube links", "learn more", "reference links",
            "study material", "more information", "show resources",
        ]),
        "regen_doc": (
            wants_doc
            or (has_n_records and bool(toks & action_toks))
            or "downloadable" in ql
            or any(w in ql for w in [
                "regenerate", "shorten doc", "rewrite doc", "create document",
                "generate document", "make document", "regenerate document",
                "updated document", "updated doc", "download document",
                "download doc",
            ])
        ),
    }

def llm_call(tpl, vars_, temp=0.3):
    """
    Call the LLM with automatic Groq → Gemini fallback.

    Previously: a single un-guarded invoke() — any Groq 429 propagated
    straight to the caller's except block and dumped raw document text.

    Now:
    • Try Groq first.
    • On 429 / rate-limit → immediately switch to Gemini (no wait).
    • Any other Groq error → re-raise so the caller knows what failed.
    • If Groq is unavailable → go straight to Gemini.
    """
    from langchain_core.prompts import PromptTemplate
    prompt = PromptTemplate(input_variables=list(vars_.keys()), template=tpl)

    def _extract(r):
        return (r.content if hasattr(r, "content") else str(r)).strip()

    groq = _llm_groq()
    if groq:
        try:
            return _extract((prompt | groq).invoke(vars_))
        except Exception as e:
            err = str(e)
            if any(tok in err for tok in ("429", "rate_limit", "rate limit",
                                          "RESOURCE_EXHAUSTED", "quota")):
                logger.warning("Groq rate-limited → switching to Gemini")
                # fall through to Gemini below
            else:
                raise   # real error — let the caller handle it

    # Gemini fallback (or primary if no Groq key)
    gemini = _llm_gemini()
    return _extract((prompt | gemini).invoke(vars_))




def direct_summarize(doc_text, action, mode="detailed", chat_history=""):
    if not doc_text or not doc_text.strip():
        return "No document content. Please upload a document first."

    mode_instructions = {
        "detailed":  "Thorough, well-structured response with all relevant details.",
        "quick":     "2-4 sentences. Direct and specific.",
        "bullet":    "Clear bullet points with specific facts.",
        "beginner":  "Simple everyday language. No jargon.",
        "executive": "Key facts only. Brief and actionable.",
        "table":     "Markdown table where appropriate.",
    }

    action_lower = action.lower()
    import re as _re

    if any(p in action_lower for p in ["five-year","5-year","child","kid","simply","beginner"]):
        mode = "beginner"

    length_match = _re.search(
        r'\bin\s+(\d+)\s+(lines?|sentences?|words?|points?|bullets?)\b', action_lower)
    if length_match:
        n, unit     = length_match.group(1), length_match.group(2)
        length_rule = f"\nCRITICAL: Answer in EXACTLY {n} {unit}. No headers. No extra text."
        mode        = "quick"
    elif any(w in action_lower for w in ["briefly","short","concise","quick"]):
        length_rule = "\nKeep answer SHORT — 2-3 sentences max. No headers."
        mode        = "quick"
    else:
        length_rule = ""

    mode_text     = mode_instructions.get(mode, mode_instructions["detailed"])
    history_block = f"CONVERSATION HISTORY:\n{chat_history}\n\n" if chat_history else ""

    tpl = """You are DocuMind AI — a highly intelligent, domain-agnostic document analyst.
Works with ANY document: books, surveys, reports, legal, financial, scientific, etc.

DOCUMENT CONTENT:
{doc_text}

{history_block}TASK: {action}
{length_rule}
STYLE: {mode_text}

RULES:
1. Read ALL content carefully before answering
2. Answer EXACTLY what is asked — no unrelated sections
3. Adapt to the document type intelligently:
   - Book/article: key ideas, frameworks, insights, actionable points
   - Survey/questionnaire: what it measures, patterns, totals, findings
   - Report: key data, conclusions, implications
   - Legal/financial: terms, figures, obligations
4. For scoring/questionnaire data: identify the scoring system from the document itself,
   compute totals dynamically, classify based on ranges in the document
5. Never add Introduction or Conclusion headers unless asked
6. Never reproduce raw data rows verbatim — interpret and explain
7. Use actual names, numbers, facts from the document
8. If not in document, say so — never hallucinate
9. For follow-ups: use conversation history, never repeat yourself

ANSWER:"""

    try:
        ans = llm_call(tpl, {
            "doc_text":      doc_text[:8000],
            "action":        action,
            "mode_text":     mode_text,
            "length_rule":   length_rule,
            "history_block": history_block,
        })
        if not ans or ans.strip().lower() in {"none", "result: none", ""}:
            raise ValueError("empty")
        return ans
    except Exception as e:
        logger.error(f"direct_summarize: {e}")
        return "The AI service is temporarily busy. Please try again in a moment."



def regen_document(doc_text, instruction):
    tpl = """You are a document transformation engine.
Output ONLY the transformed document content — nothing else.

ORIGINAL DOCUMENT:
{doc_text}

INSTRUCTION: {instruction}

ABSOLUTE RULES:
- Output the transformed document data only
- NO introduction, NO explanation, NO conclusion, NO preamble
- NO section headers like "Selected Records" or "Generated Document"
- NO sentences like "Here are the first N records"
- If instruction says "N records" → output exactly N records in the same format
- Keep all original field names and column headers exactly as they appear

OUTPUT:"""
    try:
        content = llm_call(tpl, {"doc_text": doc_text[:10000], "instruction": instruction}, temp=0.0)
    except Exception as e:
        logger.error(f"regen_document: {e}")
        content = doc_text[:3000]
    return {
        "docx": gen_file(content, "docx"),
        "pdf":  gen_file(content, "pdf"),
        "csv":  gen_file(content, "csv"),
        "txt":  gen_file(content, "txt"),
    }, content


def _markdown_table_from_df(df):
    """Build a markdown table without requiring optional tabulate dependency."""
    headers = [str(c) for c in df.columns]
    rows = df.fillna("").astype(str).values.tolist()
    out = [
        "| " + " | ".join(headers) + " |",
        "| " + " | ".join(["---"] * len(headers)) + " |",
    ]
    for row in rows:
        out.append("| " + " | ".join(str(v).replace("\n", " ").replace("|", "/") for v in row) + " |")
    return "\n".join(out)

def _extract_text_records(doc_text: str, n: int) -> list:
    import re as _re
    blocks = _re.split(r"\n{2,}", doc_text.strip())
    blocks = [b.strip() for b in blocks if b.strip()]
    if len(blocks) >= n:
        return blocks[:n]

    lines = [l.strip() for l in doc_text.split("\n") if l.strip()]
    sep_pat = _re.compile(r"^(?:\d+[\.\)]\s+.{3,}|[-=]{3,}|[A-Z][a-z]+(?:\s+[A-Z]?[a-z]+){0,4})$")
    starts = [i for i, l in enumerate(lines) if sep_pat.match(l)]
    if len(starts) >= n:
        records = []
        for idx, start in enumerate(starts[:n]):
            end = starts[idx + 1] if idx + 1 < len(starts) else len(lines)
            records.append("\n".join(lines[start:end]))
        return records

    chunk_size = max(len(lines) // max(n, 1), 1)
    return ["\n".join(lines[i * chunk_size:(i + 1) * chunk_size]) for i in range(min(n, len(lines)))]


def _split_records(doc_text: str, n: int) -> list:
    """Split text into n individual records using blank lines or headings."""
    # Blank-line separated blocks (most formats)
    blocks = [b.strip() for b in re.split(r"\n{2,}", doc_text.strip()) if b.strip()]
    if len(blocks) >= n:
        return blocks[:n]
 
    # Line-level pattern detection
    lines   = [l.strip() for l in doc_text.split("\n") if l.strip()]
    sep_pat = re.compile(
        r"^(?:\d+[\.\)]\s+.{3,}|[-=]{3,}|[A-Z][a-z]+(?:\s+[A-Z]?[a-z]+){0,4})$"
    )
    starts = [i for i, l in enumerate(lines) if sep_pat.match(l)]
    if len(starts) >= n:
        records = []
        for idx, start in enumerate(starts[:n]):
            end = starts[idx + 1] if idx + 1 < len(starts) else len(lines)
            records.append("\n".join(lines[start:end]))
        return records
 
    # Even split as last resort
    chunk = max(len(lines) // max(n, 1), 1)
    return ["\n".join(lines[i * chunk:(i + 1) * chunk]) for i in range(min(n, len(lines)))]


def _extract_score(val) -> float:
    """
    Extract the leading numeric value from any scored cell.
    Handles: '3', '2.5', '1 Some label text', 'Score: 4'
    Fully generic — no assumptions about the scoring system.
    """
    try:
        s = str(val).strip()
        import re as _re
        m = _re.match(r'^([0-9]+(?:\.[0-9]+)?)', s)
        return float(m.group(1)) if m else 0.0
    except Exception:
        return 0.0


def _smart_transform(df, instruction: str):
    """
    Transform a DataFrame based on a natural language instruction.
    Fully generic — derives all column names from the uploaded file at runtime.
    Calls assert_no_domain_leakage() to prevent domain terms creeping back in.

    Handles:
      - Group score computation: "total X score" → detects X-columns dynamically
      - Column filtering: "only include col1, col2" → matches actual column names
      - Record limiting: "first N records"
    """
    from utils.validator import assert_no_domain_leakage
    import re as _re

    assert_no_domain_leakage(instruction, context="_smart_transform instruction")

    ql     = instruction.lower()
    result = df.copy()

    # ── Dynamic group-score computation ───────────────────────────────────────
    # Detects "total <keyword> score" and computes from matching real columns.
    score_req = _re.search(
        r'total\s+(\w+)\s+score|(\w+)\s+(?:total|sum)\s+score|'
        r'compute\s+(\w+)\s+score|sum\s+(?:of\s+)?(\w+)',
        ql
    )
    if score_req:
        keyword     = next(g for g in score_req.groups() if g)
        scored_cols = [
            c for c in df.columns
            if keyword.lower() in c.lower() and
            any(w in c.lower() for w in ["question", "q", "item", "score"])
        ]
        if not scored_cols:
            scored_cols = [c for c in df.columns if keyword.lower() in c.lower()]
        if scored_cols:
            col_label      = f"Total {keyword.upper()} Score"
            result[col_label] = df[scored_cols].applymap(_extract_score).sum(axis=1)

    # ── Dynamic column filtering ──────────────────────────────────────────────
    # Builds wanted-column list entirely from what the user said and what
    # columns actually exist — zero hardcoded column name mappings.
    _wants_filter = any(p in ql for p in [
        "only include", "remove column", "remove the column",
        "keep only", "just include", "filter column",
        "select only", "show only",
    ])

    if _wants_filter:
        col_lower = {c.lower().strip(): c for c in result.columns}
        wanted    = []

        # Exact match: column name mentioned verbatim in instruction
        for col_key, col_real in col_lower.items():
            if col_key in ql:
                wanted.append(col_real)

        # Partial match: any significant word from a column name appears in the instruction
        if not wanted:
            for col_key, col_real in col_lower.items():
                words = [w for w in col_key.split() if len(w) > 2]
                if words and any(w in ql for w in words):
                    wanted.append(col_real)

        # Deduplicate while preserving order
        seen   = set()
        wanted = [c for c in wanted if not (c in seen or seen.add(c))]

        if wanted:
            result = result[[c for c in wanted if c in result.columns]]

    # ── Record limiting ───────────────────────────────────────────────────────
    limit = _re.search(
        r"\b(?:first|only|top)?\s*(\d+)\s+(?:records?|rows?|entries|people|respondents?)\b",
        ql,
    )
    if limit:
        result = result.head(int(limit.group(1)))

    return result


def build_generated_document(doc_text, instruction, file_path=""):
    """
    Generates properly formatted downloadable documents.
    CSV/Excel → pandas direct output (proper columns).
    Text docs → record extraction.
    Other transformations → LLM.
    """
    import io

    # Detect record limit
    limit_match = re.search(
        r"\b(?:only\s+|first\s+|top\s+|just\s+)?(\d+)\s+"
        r"(?:records?|rows?|entries|entry|people|persons?|respondents?|items?)\b",
        instruction.lower(),
    )
    record_limit = int(limit_match.group(1)) if limit_match else None
    ext          = os.path.splitext(file_path or "")[1].lower()

    # ── CSV / Excel → generate directly from pandas ───────────────────────────
    if file_path and ext in {".csv", ".xlsx", ".xls"}:
        try:
            import pandas as pd
            df      = pd.read_csv(file_path) if ext == ".csv" else pd.read_excel(file_path)
            limited = df.head(record_limit) if record_limit else df
            n, total = len(limited), len(df)
            label   = f"First {n} of {total} Records" if record_limit else f"All {total} Records"

            # ── Excel: proper columns via openpyxl ────────────────────────────
            xlsx_buf = io.BytesIO()
            with pd.ExcelWriter(xlsx_buf, engine="openpyxl") as writer:
                limited.to_excel(writer, index=False, sheet_name=label[:31])
            xlsx_bytes = xlsx_buf.getvalue()

            # ── CSV: raw pandas output (proper commas, no wrapper) ────────────
            csv_bytes = limited.to_csv(index=False).encode("utf-8")

            # ── DOCX / PDF: formatted markdown table ──────────────────────────
            try:
                table   = limited.to_markdown(index=False)
            except Exception:
                table   = limited.to_string(index=False)
            content = f"# {label}\n\n{table}"
            docx_bytes = gen_file(content, "docx")
            pdf_bytes  = gen_file(content, "pdf")

            return {
                "docx": docx_bytes,
                "pdf":  pdf_bytes,
                "csv":  csv_bytes,
                "txt":  csv_bytes,
            }, content

        except Exception as e:
            logger.error(f"CSV/Excel document generation failed: {e}")

    # ── Text document with record limit ───────────────────────────────────────
    if record_limit and doc_text:
        try:
            records = _split_records(doc_text, record_limit)
            if records:
                content = "\n\n---\n\n".join(records)
                return {
                    "docx": gen_file(content, "docx"),
                    "pdf":  gen_file(content, "pdf"),
                    "csv":  gen_file(content, "csv"),
                    "txt":  gen_file(content, "txt"),
                }, content
        except Exception as e:
            logger.error(f"Text record split failed: {e}")

    # ── LLM fallback for all other transformations ────────────────────────────
    return regen_document(doc_text, instruction)





def _make_fig(labels, values, title, ctype):
    try:
        import plotly.graph_objects as go
        C = ["#6366f1","#06b6d4","#8b5cf6","#10b981","#f59e0b","#ef4444","#ec4899","#14b8a6"]
        L = dict(paper_bgcolor="rgba(0,0,0,0)", plot_bgcolor="rgba(0,0,0,0)",
                 font=dict(color="#94a3b8", size=13),
                 title=dict(text=title, font=dict(color="#e2e8f0", size=16), x=0.5),
                 margin=dict(t=60, b=50, l=50, r=30), height=380, showlegend=(ctype == "pie"))
        if ctype == "pie":
            fig = go.Figure(go.Pie(labels=labels, values=values, hole=0.4,
                                   marker=dict(colors=C[:len(labels)]), textinfo="percent+label"))
        elif ctype == "line":
            fig = go.Figure(go.Scatter(x=labels, y=values, mode="lines+markers",
                                       line=dict(color="#6366f1", width=3),
                                       marker=dict(size=8, color="#06b6d4"),
                                       fill="tozeroy", fillcolor="rgba(99,102,241,0.08)"))
            L.update(xaxis=dict(gridcolor="rgba(255,255,255,0.05)"), yaxis=dict(gridcolor="rgba(255,255,255,0.05)"))
        else:
            fig = go.Figure(go.Bar(x=labels, y=values, marker=dict(color=C[:len(labels)])))
            L.update(xaxis=dict(gridcolor="rgba(255,255,255,0.03)"), yaxis=dict(gridcolor="rgba(255,255,255,0.05)"), bargap=0.3)
        fig.update_layout(**L)
        return fig
    except Exception:
        return None

def gen_charts(doc_text, question):
    """Generate charts from document data. Returns list of (title, fig) tuples."""
    tpl = """You are a data visualization expert. Extract numerical data from this document to create meaningful charts.

Task: {question}

Document/Data:
{text}

CRITICAL INSTRUCTIONS:
1. Find ACTUAL numerical values in the document
2. For spreadsheet/CSV data: count occurrences, calculate averages, find distributions
3. For text documents: extract any mentioned numbers, statistics, or counts
4. Create 1-3 charts showing the most meaningful insights
5. Return ONLY a valid JSON array - NO markdown, NO explanation, NO code blocks

Required format:
[{{"labels":["Category A","Category B","Category C"],"values":[25,40,35],"title":"Descriptive Chart Title","type":"bar"}}]

Rules:
- labels and values MUST have the same length
- values must be numbers only (no strings)
- type must be: bar, pie, or line
- If document has no numerical data, estimate distributions from categories/counts
- Generate at least 1 chart always"""
    try:
        raw = llm_call(tpl, {"text": doc_text[:4000], "question": question}, temp=0)
        raw = re.sub(r"```[a-z]*", "", raw).replace("```", "").strip()
        s = raw.find("[")
        e = raw.rfind("]") + 1
        if s < 0 or e <= s:
            s = raw.find("{")
            e = raw.rfind("}") + 1
            if s >= 0 and e > s:
                raw = "[" + raw[s:e] + "]"
                s, e = 0, len(raw)
            else:
                return []
        figs = []
        parsed = json.loads(raw[s:e])
        if isinstance(parsed, dict):
            parsed = [parsed]
        for cd in parsed[:3]:
            lb = cd.get("labels", [])
            vl = cd.get("values", [])
            title = cd.get("title", "Chart")
            ctype = cd.get("type", "bar")
            if not lb or not vl:
                continue
            if len(lb) != len(vl):
                min_len = min(len(lb), len(vl))
                lb, vl = lb[:min_len], vl[:min_len]
            try:
                vl_float = [float(str(v).replace(",", "").replace("%", "")) for v in vl]
            except Exception:
                continue
            f = _make_fig(lb, vl_float, title, ctype)
            if f:
                figs.append((title, f))
        return figs
    except Exception as e:
        logger.error(f"gen_charts: {e}")
        return []




def gen_quiz(doc_text):
    """
    Generate 5 MCQ questions that test real understanding of the document.
    Samples full document. Options shuffled so correct answer is never always A.
    """
    import random

    length = len(doc_text)
    chunk = 1000

    # Sample intelligently from 4 parts of the document
    sections = [
        doc_text[:chunk],
        doc_text[length // 3 : length // 3 + chunk],
        doc_text[2 * length // 3 : 2 * length // 3 + chunk],
        doc_text[max(0, length - chunk):]
    ]
    sampled = "\n\n---\n\n".join(s.strip() for s in sections if s.strip())

    tpl = """You are creating a quiz to test whether someone has carefully read this document.

Document content (sampled from throughout):
{text}

Generate exactly 5 multiple choice questions. Each question must:
- Test REAL understanding — not obvious or trivial facts
- Reference a SPECIFIC detail (name, number, finding, category, result)
  that is actually in the document above
- Have 4 options where the wrong options are PLAUSIBLE (not obviously wrong)
- Cover DIFFERENT sections/topics of the document
- Range from factual recall to analytical reasoning

FORBIDDEN question types:
- "What is the title of the document?"
- "Who wrote this document?"
- Questions answerable without reading the document
- Trick questions
- Questions where all wrong answers are obviously wrong

Use EXACTLY this format (blank line between questions):
Q: [specific, meaningful question — 10-20 words]
CORRECT: [the correct answer from the document]
WRONG1: [plausible but incorrect option]
WRONG2: [plausible but incorrect option]
WRONG3: [plausible but incorrect option]

5 questions:"""

    try:
        raw = llm_call(tpl, {"text": sampled}, temp=0.3)
        questions = []
        labels = ["A", "B", "C", "D"]

        for block in re.split(r"\n\s*\n", raw.strip()):
            parsed = {}
            for ln in [l.strip() for l in block.strip().split("\n") if l.strip()]:
                if ":" in ln:
                    key, val = ln.split(":", 1)
                    parsed[key.strip().upper()] = val.strip()

            q_text = parsed.get("Q", "")
            correct_text = parsed.get("CORRECT", "")
            wrongs = [
                parsed.get("WRONG1", ""),
                parsed.get("WRONG2", ""),
                parsed.get("WRONG3", "")
            ]
            wrongs = [w for w in wrongs if w and w.lower() != correct_text.lower()]

            # Skip trivial or forbidden questions
            if not q_text or not correct_text or len(wrongs) < 2:
                continue
            q_lower = q_text.lower()
            forbidden_q = [
                "what is the title", "who wrote", "who is the author",
                "what is the name of the document", "what type of document"
            ]
            if any(f in q_lower for f in forbidden_q):
                continue

            # Build and shuffle options
            all_opts = [correct_text] + wrongs[:3]
            while len(all_opts) < 4:
                all_opts.append("Not mentioned in the document")
            random.shuffle(all_opts)

            try:
                correct_label = labels[all_opts.index(correct_text)]
            except ValueError:
                continue

            options = [(labels[i], all_opts[i]) for i in range(len(all_opts))]
            questions.append({
                "question": q_text,
                "options": options,
                "answer": correct_label
            })

        return questions[:5]

    except Exception as e:
        logger.error(f"gen_quiz: {e}")
        return []

def _render_quiz_legacy(quiz, ks):
    if not quiz:
        st.warning("Could not generate quiz. Try asking again.")
        return
    st.markdown(
        "<div style='background:linear-gradient(135deg,rgba(139,92,246,0.1),rgba(99,102,241,0.06));"
        "border:1px solid rgba(139,92,246,0.25);border-radius:16px;padding:18px 22px;margin:16px 0 12px;'>"
        "<span style='font-size:0.62rem;font-weight:800;color:rgba(139,92,246,0.9);"
        "text-transform:uppercase;letter-spacing:3px;'>📝 Interactive Quiz</span>"
        "<p style='color:#64748b;font-size:0.78rem;margin:5px 0 0;'>Select your answer and click Check.</p></div>",
        unsafe_allow_html=True)
    for i, q in enumerate(quiz, 1):
        st.markdown(
            f"<div style='background:rgba(8,10,22,0.7);border:1px solid rgba(139,92,246,0.18);"
            f"border-radius:12px;padding:16px 20px;margin:10px 0;'>"
            f"<span style='color:rgba(139,92,246,0.5);font-size:0.6rem;font-weight:800;"
            f"text-transform:uppercase;letter-spacing:2px;'>Question {i}</span>"
            f"<p style='color:#e2e8f0;font-size:0.9rem;font-weight:600;margin:8px 0 0;'>{q['question']}</p></div>",
            unsafe_allow_html=True)
        if q.get("options"):

            import hashlib as _hl
            qhash = _hl.md5(q.get("question","").encode()).hexdigest()[:6]
            chosen = st.radio(f"q{i}", [f"{o[0]}.  {o[1]}" for o in q["options"]],
                              key=f"quiz_r_{i}_{ks}_{qhash}", label_visibility="collapsed")
            if st.button(f"Check Q{i}", key=f"quiz_b_{i}_{ks}_{qhash}"):
                correct = q.get("answer", "A")
                if chosen and chosen[0] == correct:
                    st.success("Correct!")
                else:
                    ct = next((o[1] for o in q["options"] if o[0] == correct), correct)
                    st.error(f"Correct: {correct}.  {ct}")




def render_quiz(quiz, ks):
    if not quiz:
        st.warning("Could not generate quiz. Try asking again.")
        return
    import hashlib as _hl
    st.markdown(
        "<div style='background:linear-gradient(135deg,rgba(139,92,246,0.1),rgba(99,102,241,0.06));"
        "border:1px solid rgba(139,92,246,0.25);border-radius:16px;padding:18px 22px;margin:16px 0 12px;'>"
        "<span style='font-size:0.62rem;font-weight:800;color:rgba(139,92,246,0.9);"
        "text-transform:uppercase;letter-spacing:3px;'>Interactive Quiz</span>"
        "<p style='color:#64748b;font-size:0.78rem;margin:6px 0 0;'>Select your answer then click Check. No option is pre-selected.</p></div>",
        unsafe_allow_html=True)
    score_key = f"quiz_score_{ks}"
    if score_key not in st.session_state:
        st.session_state[score_key] = {}
    for i, q in enumerate(quiz, 1):
        qhash = _hl.md5(q.get("question", str(i)).encode()).hexdigest()[:8]
        st.markdown(
            f"<div style='background:rgba(8,10,22,0.75);border:1px solid rgba(139,92,246,0.18);"
            f"border-radius:12px;padding:16px 20px;margin:12px 0;'>"
            f"<span style='color:rgba(139,92,246,0.55);font-size:0.59rem;font-weight:800;"
            f"text-transform:uppercase;letter-spacing:2px;'>Question {i} of {len(quiz)}</span>"
            f"<p style='color:#e2e8f0;font-size:0.88rem;font-weight:600;margin:8px 0 0;line-height:1.5;'>"
            f"{q['question']}</p></div>",
            unsafe_allow_html=True)
        if q.get("options"):
            radio_key = f"quiz_r_{i}_{ks}_{qhash}"
            chosen = st.radio(
                f"Select answer for question {i}:",
                [f"{o[0]}.  {o[1]}" for o in q["options"]],
                key=radio_key,
                label_visibility="collapsed",
                index=None
            )
            result_key = f"quiz_result_{i}_{ks}_{qhash}"
            if st.button(f"Check Q{i}", key=f"quiz_b_{i}_{ks}_{qhash}"):
                if chosen is None:
                    st.warning("Please select an option before checking.")
                else:
                    correct = q.get("answer", "A")
                    is_correct = chosen[0] == correct
                    st.session_state[score_key][i] = is_correct
                    st.session_state[result_key] = (is_correct, correct, q.get("options", []))
            if result_key in st.session_state:
                is_correct, correct, options = st.session_state[result_key]
                if is_correct:
                    st.success("Correct!")
                else:
                    ct = next((o[1] for o in options if o[0] == correct), correct)
                    st.error(f"Correct answer: {correct}.  {ct}")
    answered = st.session_state[score_key]
    if answered:
        correct_count = sum(1 for v in answered.values() if v)
        total_answered = len(answered)
        pct = int(correct_count / len(quiz) * 100)
        color = "#3fb950" if pct >= 70 else "#f59e0b" if pct >= 40 else "#ef4444"
        st.markdown(
            f"<div style='background:rgba(8,10,22,0.6);border:1px solid rgba(99,102,241,0.15);"
            f"border-radius:10px;padding:12px 18px;margin-top:14px;text-align:center;'>"
            f"<span style='color:{color};font-size:1rem;font-weight:700;'>"
            f"Score: {correct_count}/{len(quiz)} ({pct}%)</span>"
            f"<span style='color:#64748b;font-size:0.75rem;margin-left:12px;'>"
            f"{total_answered} of {len(quiz)} answered</span></div>",
            unsafe_allow_html=True)

def show_resources(question, answer):
    tpl = """Generate search topics for: {question}
Context: {answer}
Return ONLY this JSON (no markdown):
{{"youtube":["topic1","topic2","topic3"],"google":["phrase1","phrase2"],"scholar":["term1","term2"]}}"""
    with st.spinner("Finding resources..."):
        try:
            raw = llm_call(tpl, {"question": question, "answer": answer[:300]}, temp=0.1)
            raw = re.sub(r"```[a-z]*", "", raw).replace("```", "").strip()
            s, e = raw.find("{"), raw.rfind("}") + 1
            k = json.loads(raw[s:e]) if s >= 0 else {}
        except Exception:
            k = {}
    if not k: st.info("No resources found."); return
    def enc(t): return t.strip().replace(" ", "+").replace("&", "%26")
    st.markdown("<div style='margin:20px 0 14px;padding:12px 16px;background:linear-gradient(135deg,rgba(99,102,241,0.06),rgba(6,182,212,0.03));border:1px solid rgba(99,102,241,0.15);border-radius:12px;'><span style='font-size:0.62rem;font-weight:800;color:rgba(99,102,241,0.8);text-transform:uppercase;letter-spacing:3px;'>Learning Resources</span></div>", unsafe_allow_html=True)
    c1, c2, c3 = st.columns(3)
    with c1:
        st.markdown("<p style='font-size:0.73rem;font-weight:700;color:#fca5a5;margin-bottom:8px;'>YouTube</p>", unsafe_allow_html=True)
        for s in (k.get("youtube") or [])[:3]:
            st.markdown(f"<a href='https://www.youtube.com/results?search_query={enc(s)}' target='_blank' style='display:block;background:rgba(239,68,68,0.07);border:1px solid rgba(239,68,68,0.18);border-radius:9px;padding:9px 13px;margin:5px 0;color:#fca5a5;text-decoration:none;font-size:0.78rem;font-weight:500;'>{s}</a>", unsafe_allow_html=True)
    with c2:
        st.markdown("<p style='font-size:0.73rem;font-weight:700;color:#a5b4fc;margin-bottom:8px;'>Google</p>", unsafe_allow_html=True)
        for t in (k.get("google") or [])[:3]:
            st.markdown(f"<a href='https://www.google.com/search?q={enc(t)}' target='_blank' style='display:block;background:rgba(99,102,241,0.07);border:1px solid rgba(99,102,241,0.15);border-radius:9px;padding:9px 13px;margin:5px 0;color:#a5b4fc;text-decoration:none;font-size:0.78rem;font-weight:500;'>{t}</a>", unsafe_allow_html=True)
    with c3:
        st.markdown("<p style='font-size:0.73rem;font-weight:700;color:#67e8f9;margin-bottom:8px;'>Scholar</p>", unsafe_allow_html=True)
        for s in (k.get("scholar") or [])[:3]:
            st.markdown(f"<a href='https://scholar.google.com/scholar?q={enc(s)}' target='_blank' style='display:block;background:rgba(6,182,212,0.05);border:1px solid rgba(6,182,212,0.15);border-radius:9px;padding:9px 13px;margin:5px 0;color:#67e8f9;text-decoration:none;font-size:0.78rem;font-weight:500;'>{s}</a>", unsafe_allow_html=True)




def _parse_markdown_lines(content):
    """Parse markdown content into structured line objects for formatting."""
    lines = []
    for ln in content.split("\n"):
        raw = ln.rstrip()
        if raw.startswith("### "):
            lines.append(("h3", raw[4:].strip()))
        elif raw.startswith("## "):
            lines.append(("h2", raw[3:].strip()))
        elif raw.startswith("# "):
            lines.append(("h1", raw[2:].strip()))
        elif raw.startswith("**") and raw.endswith("**") and len(raw) > 4:
            lines.append(("bold", raw[2:-2].strip()))
        elif raw.startswith("  - ") or raw.startswith("  * "):
            lines.append(("bullet2", raw[4:].strip()))
        elif raw.startswith("- ") or raw.startswith("* "):
            lines.append(("bullet", raw[2:].strip()))
        elif "|" in raw and raw.count("|") >= 2:
            cols = [c.strip() for c in raw.split("|") if c.strip() and set(c.strip()) != {"-"}]
            if cols:
                lines.append(("table_row", cols))
        elif raw.strip() == "":
            lines.append(("blank", ""))
        else:
            lines.append(("text", raw.strip()))
    return lines


def gen_file_docx(content, title="DocuMind AI Report"):
    """Generate properly formatted DOCX with styles, bullets, tables, spacing."""
    from docx import Document
    from docx.shared import Pt, RGBColor, Inches
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.enum.table import WD_ALIGN_VERTICAL, WD_TABLE_ALIGNMENT

    doc = Document()
    for section in doc.sections:
        section.top_margin = Inches(1)
        section.bottom_margin = Inches(1)
        section.left_margin = Inches(1.2)
        section.right_margin = Inches(1.2)

    title_para = doc.add_heading(title, 0)
    title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    title_run = title_para.runs[0] if title_para.runs else title_para.add_run(title)
    title_run.font.color.rgb = RGBColor(0x1e, 0x29, 0x3b)

    import datetime
    meta = doc.add_paragraph(f"Generated by DocuMind AI - {datetime.datetime.now().strftime('%B %d, %Y')}")
    meta.alignment = WD_ALIGN_PARAGRAPH.CENTER
    meta.runs[0].font.size = Pt(9)
    meta.runs[0].font.color.rgb = RGBColor(0x64, 0x74, 0x8b)
    doc.add_paragraph()

    table_buf = []

    def flush_table():
        nonlocal table_buf
        if not table_buf:
            return
        headers = table_buf[0]
        data_rows = table_buf[1:]
        if headers:
            tbl = doc.add_table(rows=1 + len(data_rows), cols=len(headers))
            tbl.style = "Table Grid"
            tbl.alignment = WD_TABLE_ALIGNMENT.CENTER
            hdr_cells = tbl.rows[0].cells
            for ci, hdr in enumerate(headers):
                hdr_cells[ci].text = hdr
                hdr_cells[ci].vertical_alignment = WD_ALIGN_VERTICAL.CENTER
                run = hdr_cells[ci].paragraphs[0].runs[0] if hdr_cells[ci].paragraphs[0].runs else hdr_cells[ci].paragraphs[0].add_run(hdr)
                run.font.bold = True
                run.font.size = Pt(10)
                hdr_cells[ci].paragraphs[0].alignment = WD_ALIGN_PARAGRAPH.CENTER
            for ri, row in enumerate(data_rows):
                row_cells = tbl.rows[ri + 1].cells
                for ci, val in enumerate(row[:len(headers)]):
                    row_cells[ci].text = val
                    row_cells[ci].vertical_alignment = WD_ALIGN_VERTICAL.CENTER
                    if row_cells[ci].paragraphs[0].runs:
                        row_cells[ci].paragraphs[0].runs[0].font.size = Pt(10)
            doc.add_paragraph()
        table_buf = []

    for ltype, lval in _parse_markdown_lines(content):
        if ltype == "table_row":
            table_buf.append(lval)
            continue

        flush_table()
        if ltype == "h1":
            p = doc.add_heading(lval, 1)
            if p.runs:
                p.runs[0].font.color.rgb = RGBColor(0x1e, 0x29, 0x3b)
        elif ltype in ("h2", "bold"):
            p = doc.add_heading(lval, 2)
            if p.runs:
                p.runs[0].font.color.rgb = RGBColor(0x1e, 0x40, 0xaf)
        elif ltype == "h3":
            doc.add_heading(lval, 3)
        elif ltype == "bullet":
            p = doc.add_paragraph(style="List Bullet")
            run = p.add_run(lval)
            run.font.size = Pt(11)
        elif ltype == "bullet2":
            p = doc.add_paragraph(style="List Bullet 2")
            run = p.add_run(lval)
            run.font.size = Pt(10)
        elif ltype == "blank":
            doc.add_paragraph()
        elif ltype == "text" and lval:
            p = doc.add_paragraph()
            parts = re.split(r"(\*\*.*?\*\*)", lval)
            for part in parts:
                if part.startswith("**") and part.endswith("**"):
                    run = p.add_run(part[2:-2])
                    run.bold = True
                    run.font.size = Pt(11)
                elif part:
                    run = p.add_run(part)
                    run.font.size = Pt(11)
            p.paragraph_format.space_after = Pt(4)

    flush_table()
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def gen_file_pdf(content, title="DocuMind AI Report"):
    """Generate properly formatted PDF with headers, bullets, tables, fonts."""
    try:
        from html import escape
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib.units import cm
        from reportlab.lib import colors
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
        from reportlab.lib.enums import TA_CENTER

        buf = io.BytesIO()
        doc = SimpleDocTemplate(
            buf, pagesize=A4,
            leftMargin=2.5 * cm, rightMargin=2.5 * cm,
            topMargin=2.5 * cm, bottomMargin=2.5 * cm
        )

        base = getSampleStyleSheet()
        styles = {
            "title": ParagraphStyle("DM_Title", parent=base["Title"], fontSize=22,
                textColor=colors.HexColor("#1e293b"), spaceAfter=6, alignment=TA_CENTER,
                fontName="Helvetica-Bold"),
            "meta": ParagraphStyle("DM_Meta", parent=base["Normal"], fontSize=9,
                textColor=colors.HexColor("#64748b"), spaceAfter=16, alignment=TA_CENTER),
            "h1": ParagraphStyle("DM_H1", parent=base["Heading1"], fontSize=16,
                textColor=colors.HexColor("#1e293b"), spaceBefore=14, spaceAfter=6,
                fontName="Helvetica-Bold"),
            "h2": ParagraphStyle("DM_H2", parent=base["Heading2"], fontSize=13,
                textColor=colors.HexColor("#1e40af"), spaceBefore=10, spaceAfter=4,
                fontName="Helvetica-Bold"),
            "h3": ParagraphStyle("DM_H3", parent=base["Heading3"], fontSize=11,
                textColor=colors.HexColor("#374151"), spaceBefore=8, spaceAfter=3,
                fontName="Helvetica-Bold"),
            "body": ParagraphStyle("DM_Body", parent=base["Normal"], fontSize=10.5,
                textColor=colors.HexColor("#1e293b"), spaceAfter=5, leading=16),
            "bullet": ParagraphStyle("DM_Bullet", parent=base["Normal"], fontSize=10.5,
                textColor=colors.HexColor("#374151"), leftIndent=18, firstLineIndent=-12,
                spaceAfter=3, leading=15),
            "bullet2": ParagraphStyle("DM_Bullet2", parent=base["Normal"], fontSize=10,
                textColor=colors.HexColor("#4b5563"), leftIndent=36, firstLineIndent=-12,
                spaceAfter=3, leading=14),
        }

        import datetime
        story = [
            Paragraph(title, styles["title"]),
            Paragraph(f"Generated by DocuMind AI - {datetime.datetime.now().strftime('%B %d, %Y')}", styles["meta"]),
            Spacer(1, 0.3 * cm),
        ]

        def rich_text(value):
            parts = re.split(r"(\*\*.*?\*\*)", value)
            out = []
            for part in parts:
                if part.startswith("**") and part.endswith("**"):
                    out.append(f"<b>{escape(part[2:-2])}</b>")
                else:
                    out.append(escape(part))
            return "".join(out)

        table_buf = []

        def flush_table():
            nonlocal table_buf
            if not table_buf:
                return
            headers = table_buf[0]
            data_rows = table_buf[1:]
            if headers:
                tdata = []
                for row in [headers] + data_rows:
                    padded = list(row[:len(headers)]) + [""] * max(len(headers) - len(row), 0)
                    tdata.append([escape(str(c)) for c in padded])
                col_w = doc.width / max(len(headers), 1)
                t = Table(tdata, colWidths=[col_w] * len(headers), repeatRows=1)
                t.setStyle(TableStyle([
                    ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#1e40af")),
                    ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
                    ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                    ("FONTSIZE", (0, 0), (-1, 0), 10),
                    ("ALIGN", (0, 0), (-1, -1), "LEFT"),
                    ("FONTNAME", (0, 1), (-1, -1), "Helvetica"),
                    ("FONTSIZE", (0, 1), (-1, -1), 9.5),
                    ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.HexColor("#f8fafc"), colors.white]),
                    ("GRID", (0, 0), (-1, -1), 0.5, colors.HexColor("#cbd5e1")),
                    ("TOPPADDING", (0, 0), (-1, -1), 5),
                    ("BOTTOMPADDING", (0, 0), (-1, -1), 5),
                    ("LEFTPADDING", (0, 0), (-1, -1), 8),
                    ("RIGHTPADDING", (0, 0), (-1, -1), 8),
                    ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
                ]))
                story.append(t)
                story.append(Spacer(1, 0.3 * cm))
            table_buf = []

        for ltype, lval in _parse_markdown_lines(content):
            if ltype == "table_row":
                table_buf.append(lval)
                continue

            flush_table()
            if ltype == "h1":
                story.append(Paragraph(escape(lval), styles["h1"]))
            elif ltype in ("h2", "bold"):
                story.append(Paragraph(escape(lval), styles["h2"]))
            elif ltype == "h3":
                story.append(Paragraph(escape(lval), styles["h3"]))
            elif ltype == "bullet":
                story.append(Paragraph(f"- {rich_text(lval)}", styles["bullet"]))
            elif ltype == "bullet2":
                story.append(Paragraph(f"  - {rich_text(lval)}", styles["bullet2"]))
            elif ltype == "blank":
                story.append(Spacer(1, 0.15 * cm))
            elif ltype == "text" and lval:
                story.append(Paragraph(rich_text(lval), styles["body"]))

        flush_table()
        doc.build(story)
        return buf.getvalue()
    except Exception as e:
        logger.error(f"PDF generation error: {e}")
        return content.encode("utf-8")


def gen_file_xlsx(content, title="DocuMind AI Report"):
    """Generate properly formatted Excel with styled headers, alignment, column widths."""
    try:
        from openpyxl import Workbook
        from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
        from openpyxl.utils import get_column_letter

        wb = Workbook()
        ws = wb.active
        ws.title = "DocuMind Report"

        import datetime
        ws.merge_cells("A1:F1")
        ws["A1"] = title
        ws["A1"].font = Font(name="Calibri", size=16, bold=True, color="1E293B")
        ws["A1"].alignment = Alignment(horizontal="center", vertical="center")
        ws["A1"].fill = PatternFill("solid", fgColor="EFF6FF")
        ws.row_dimensions[1].height = 32

        ws.merge_cells("A2:F2")
        ws["A2"] = f"Generated by DocuMind AI - {datetime.datetime.now().strftime('%B %d, %Y')}"
        ws["A2"].font = Font(name="Calibri", size=10, color="64748B")
        ws["A2"].alignment = Alignment(horizontal="center", vertical="center")
        ws.row_dimensions[2].height = 20
        ws.row_dimensions[3].height = 8

        table_headers = None
        table_data = []
        lines = _parse_markdown_lines(content)
        for ltype, lval in lines:
            if ltype == "table_row":
                if table_headers is None:
                    table_headers = lval
                else:
                    table_data.append(lval)

        header_fill = PatternFill("solid", fgColor="1E40AF")
        alt_fill = PatternFill("solid", fgColor="F1F5F9")
        white_fill = PatternFill("solid", fgColor="FFFFFF")
        thin_border = Border(
            left=Side(style="thin", color="CBD5E1"),
            right=Side(style="thin", color="CBD5E1"),
            top=Side(style="thin", color="CBD5E1"),
            bottom=Side(style="thin", color="CBD5E1"),
        )

        if table_headers:
            header_row = 4
            ws.row_dimensions[header_row].height = 22
            for ci, hdr in enumerate(table_headers, 1):
                cell = ws.cell(row=header_row, column=ci, value=hdr)
                cell.font = Font(name="Calibri", size=11, bold=True, color="FFFFFF")
                cell.fill = header_fill
                cell.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)
                cell.border = thin_border

            for ri, row in enumerate(table_data, header_row + 1):
                ws.row_dimensions[ri].height = 18
                fill = alt_fill if ri % 2 == 0 else white_fill
                for ci, val in enumerate(row[:len(table_headers)], 1):
                    cell = ws.cell(row=ri, column=ci, value=val)
                    cell.font = Font(name="Calibri", size=10, color="1E293B")
                    cell.fill = fill
                    cell.alignment = Alignment(horizontal="left", vertical="center", wrap_text=True)
                    cell.border = thin_border

            for ci in range(1, len(table_headers) + 1):
                max_len = max(
                    len(str(ws.cell(row=r, column=ci).value or ""))
                    for r in range(header_row, header_row + 1 + len(table_data))
                )
                ws.column_dimensions[get_column_letter(ci)].width = min(max(max_len + 4, 12), 40)
        else:
            current_row = 4
            for ltype, lval in lines:
                if not lval and ltype == "blank":
                    ws.row_dimensions[current_row].height = 8
                    current_row += 1
                    continue
                if not lval:
                    continue

                cell = ws.cell(row=current_row, column=1, value=lval)
                ws.merge_cells(f"A{current_row}:F{current_row}")
                if ltype in ("h1", "h2", "h3", "bold"):
                    cell.font = Font(name="Calibri", size=12 if ltype == "h1" else 11, bold=True, color="1E40AF")
                    cell.fill = PatternFill("solid", fgColor="EFF6FF")
                    ws.row_dimensions[current_row].height = 20
                elif ltype == "bullet":
                    cell.value = f"  -  {lval}"
                    cell.font = Font(name="Calibri", size=10, color="374151")
                    ws.row_dimensions[current_row].height = 16
                else:
                    cell.font = Font(name="Calibri", size=10, color="1E293B")
                    ws.row_dimensions[current_row].height = 16
                cell.alignment = Alignment(horizontal="left", vertical="center", wrap_text=True)
                current_row += 1
            ws.column_dimensions["A"].width = 80

        buf = io.BytesIO()
        wb.save(buf)
        return buf.getvalue()
    except Exception as e:
        logger.error(f"XLSX generation error: {e}")
        import pandas as pd
        df = pd.DataFrame({"Content": [l.strip() for l in content.split("\n") if l.strip()]})
        buf = io.BytesIO()
        df.to_excel(buf, index=False, engine="openpyxl")
        return buf.getvalue()


def gen_file_csv(content):
    """Generate clean CSV; detects table structure first."""
    import pandas as pd
    lines = _parse_markdown_lines(content)
    table_headers = None
    table_data = []
    for ltype, lval in lines:
        if ltype == "table_row":
            if table_headers is None:
                table_headers = lval
            else:
                table_data.append(lval)
    if table_headers and table_data:
        try:
            df = pd.DataFrame(table_data, columns=table_headers)
            return df.to_csv(index=False).encode("utf-8")
        except Exception:
            pass
    df = pd.DataFrame({"Content": [l.strip() for l in content.split("\n") if l.strip()]})
    return df.to_csv(index=False).encode("utf-8")


def gen_file(content, fmt):
    """Unified file generation router; all formats properly formatted."""
    if fmt == "txt":
        return content.encode("utf-8")
    if fmt == "docx":
        return gen_file_docx(content)
    if fmt == "pdf":
        return gen_file_pdf(content)
    if fmt == "csv":
        return gen_file_csv(content)
    if fmt == "xlsx":
        return gen_file_xlsx(content)
    return content.encode("utf-8")


def gen_excel_from_doc(doc_text, question):
    """Extract structured tabular data from document and write to formatted Excel."""
    tpl = """Extract ALL tabular/structured data to answer: {question}
Document: {text}
Return ONLY valid JSON array of objects, no markdown:
[{{"Column1":"Value1","Column2":"Value2"}}]
Include as many rows as are in the document. Keep original column names."""
    try:
        import pandas as pd
        raw = llm_call(tpl, {"text": doc_text[:6000], "question": question}, temp=0)
        raw = re.sub(r"```[a-z]*", "", raw).replace("```", "").strip()
        s, e = raw.find("["), raw.rfind("]") + 1
        if s >= 0 and e > s:
            data = json.loads(raw[s:e])
            if data and len(data) > 0:
                df = pd.DataFrame(data)
                md_table = "| " + " | ".join(df.columns) + " |\n"
                md_table += "| " + " | ".join(["---"] * len(df.columns)) + " |\n"
                for _, row in df.iterrows():
                    md_table += "| " + " | ".join([str(v) for v in row.values]) + " |\n"
                return gen_file_xlsx(md_table, title="DocuMind AI - Data Export")
    except Exception as e: logger.error(f"gen_excel: {e}")
    return gen_file_xlsx(doc_text[:8000], title="DocuMind AI - Document Export")


def gen_pptx(content, title="DocuMind AI"):
    """Generate styled PowerPoint presentation."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        from pptx.enum.shapes import MSO_SHAPE

        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(0x04, 0x06, 0x0f)

        tx_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.33), Inches(1.5))
        tf = tx_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.alignment = PP_ALIGN.CENTER
        run = p.runs[0]
        run.font.size = Pt(44)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0xF1, 0xF5, 0xF9)

        sub_box = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(11.33), Inches(0.5))
        sp = sub_box.text_frame.paragraphs[0]
        sp.text = "Generated by DocuMind AI"
        sp.alignment = PP_ALIGN.CENTER
        srun = sp.runs[0]
        srun.font.size = Pt(18)
        srun.font.color.rgb = RGBColor(0x63, 0x66, 0xF1)

        sections = []
        current_sec = {"title": "", "bullets": []}
        for ltype, lval in _parse_markdown_lines(content):
            if ltype in ("h1", "h2", "h3", "bold") and lval:
                if current_sec["bullets"] or current_sec["title"]:
                    sections.append(current_sec)
                current_sec = {"title": lval, "bullets": []}
            elif ltype in ("bullet", "bullet2") and lval:
                current_sec["bullets"].append(lval)
            elif ltype == "text" and lval:
                current_sec["bullets"].append(lval)
        if current_sec["title"] or current_sec["bullets"]:
            sections.append(current_sec)

        for sec in sections[:15]:
            sl2 = prs.slides.add_slide(prs.slide_layouts[6])
            sl2.background.fill.solid()
            sl2.background.fill.fore_color.rgb = RGBColor(0x04, 0x06, 0x0f)

            title_box = sl2.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.8))
            tp = title_box.text_frame.paragraphs[0]
            tp.text = sec["title"][:80] if sec["title"] else "Details"
            tp.alignment = PP_ALIGN.LEFT
            trun = tp.runs[0]
            trun.font.size = Pt(28)
            trun.font.bold = True
            trun.font.color.rgb = RGBColor(0xA5, 0xB4, 0xFC)

            line = sl2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.2), Inches(12.33), Emu(18000))
            line.fill.solid()
            line.fill.fore_color.rgb = RGBColor(0x63, 0x66, 0xF1)
            line.line.fill.background()

            content_box = sl2.shapes.add_textbox(Inches(0.6), Inches(1.4), Inches(12.1), Inches(5.5))
            ctf = content_box.text_frame
            ctf.word_wrap = True
            for bi, bullet in enumerate(sec["bullets"][:10]):
                p = ctf.paragraphs[0] if bi == 0 else ctf.add_paragraph()
                p.text = f"-  {bullet}"
                p.alignment = PP_ALIGN.LEFT
                if p.runs:
                    p.runs[0].font.size = Pt(18)
                    p.runs[0].font.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
                p.space_after = Pt(6)

        buf = io.BytesIO()
        prs.save(buf)
        return buf.getvalue()
    except Exception as e: logger.error(f"gen_pptx: {e}"); return None




def render_tts(text, ks=""):
    clean = (text[:500].replace("\\","").replace("`","").replace('"',"")
             .replace("'","").replace("\n"," ").replace("\r","")
             .replace("<","").replace(">","").replace("&","and"))
    fp, fs = f"play_{ks}", f"stop_{ks}"
    st.iframe(
        f"<script>\nfunction {fp}(){{\n  window.speechSynthesis.cancel();\n"
        f"  var u=new SpeechSynthesisUtterance('{clean}');\n"
        f"  u.rate=0.92;u.pitch=1.0;u.volume=1.0;\n  window.speechSynthesis.speak(u);\n}}\n"
        f"function {fs}(){{window.speechSynthesis.cancel();}}\n</script>\n"
        f"<div style='display:flex;gap:8px;margin-top:8px;'>"
        f"<button onclick='{fp}()' style='background:rgba(99,102,241,0.12);border:1px solid rgba(99,102,241,0.28);border-radius:8px;color:#a5b4fc;padding:5px 14px;font-size:0.71rem;cursor:pointer;font-weight:600;font-family:system-ui,sans-serif;'>Read Aloud</button>"
        f"<button onclick='{fs}()' style='background:rgba(239,68,68,0.08);border:1px solid rgba(239,68,68,0.2);border-radius:8px;color:#fca5a5;padding:5px 14px;font-size:0.71rem;cursor:pointer;font-weight:600;font-family:system-ui,sans-serif;'>Stop</button>"
        f"</div>",
        height=50)




def gen_prompts(text):
    """
    Generate 6 smart, document-specific suggested questions.
    Samples intelligently from the full document.
    """
    import random
    length = len(text)

    # Sample 4 sections: start, 1/3, 2/3, end
    chunk = 900
    sections = [
        text[:chunk],
        text[length // 3 : length // 3 + chunk],
        text[2 * length // 3 : 2 * length // 3 + chunk],
        text[max(0, length - chunk):]
    ]
    sampled = "\n\n---\n\n".join(s.strip() for s in sections if s.strip())

    tpl = """You are an expert analyst who has read the following document carefully.

Document content (sampled from throughout):
{text}

Your task: Write exactly 6 questions that a professional user would ask
to get deep, useful insights from this document.

STRICT REQUIREMENTS:
- Each question MUST reference a SPECIFIC person, number, finding, term,
  category, score, percentage, or concept that actually appears above.
- Questions must require actual reading of the document to answer —
  not common knowledge.
- At least 2 questions must be ANALYTICAL (e.g. "Why did X happen?",
  "What is the significance of Y?", "How does A compare to B?").
- At least 1 question must ask about a SPECIFIC NUMBER or DATA POINT
  visible in the document.
- At least 1 question must be about a PATTERN or TREND in the data.
- FORBIDDEN question types: "What is the title?", "Who is the author?",
  "What is this document about?", "What is the purpose?",
  "What are the main topics?", any question answerable in 1 word.
- Each question: 8-15 words, ends with "?"
- One question per line. No numbering. No bullets.

6 high-quality questions:"""

    try:
        raw = llm_call(tpl, {"text": sampled}, temp=0.3)
        lines = []
        forbidden = [
            "what is this document", "what is the title",
            "who is the author", "what is the purpose",
            "what are the main", "what is the document about",
            "what is the overview", "what does this document",
            "what is the subject",
        ]
        for ln in raw.strip().split("\n"):
            ln = re.sub(r"^[\d\.\-\*\x95\u2022\)]\s*", "", ln.strip()).strip()
            if not ln or len(ln) < 15 or not ln.endswith("?"):
                continue
            lower = ln.lower()
            if any(f in lower for f in forbidden):
                continue
            lines.append(ln)
        return lines[:6] if len(lines) >= 3 else _default_prompts()
    except Exception:
        return _default_prompts()

def _default_prompts():
    return [
        "What are the key findings or arguments?",
        "What problem does this document address?",
        "What evidence supports the main claims?",
        "What actions or decisions are recommended?",
        "What risks or limitations are identified?",
        "How does this compare to alternative approaches?",
    ]




def _build_indexes(chunks, deps, file_ext=""):
    """
    CSV/Excel → BM25 only.
    Text docs → BM25 + FAISS (synchronous — used by switch_document only).
    process_doc() bypasses this and calls each builder directly.
    """
    file_ext = (file_ext or "").lower()
    if file_ext in {".csv", ".xlsx", ".xls"}:
        deps["create_bm25_index"](chunks)
        logger.debug("CSV/Excel: BM25 only (FAISS not needed)")
        return
    deps["create_bm25_index"](chunks)
    try:
        from vector_store.faiss_store import faiss_enabled
        if faiss_enabled():
            deps["create_vector_store"](chunks)   # sync — intentional for switching
    except Exception as e:
        logger.warning(f"FAISS skipped: {e}")


def process_doc(f, deps):
    import tempfile
    import hashlib

    try:
        suffix    = "." + f.name.split(".")[-1]
        raw_bytes = f.getvalue()

        with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
            tmp.write(raw_bytes)
            tmp_path = tmp.name

        ok, msg = deps["validate_file"](tmp_path)
        if not ok:
            st.error(msg)
            return False

        file_ext  = os.path.splitext(tmp_path)[1].lower()
        file_hash = hashlib.md5(raw_bytes).hexdigest()
        is_struct = file_ext in {".csv", ".xlsx", ".xls"}

        prog   = st.progress(0)
        status = st.empty()

        def step(p, lbl):
            prog.progress(p)
            status.markdown(
                f"<p style='color:#a5b4fc;font-size:0.82rem;margin:4px 0;'>⚙ {lbl}</p>",
                unsafe_allow_html=True,
            )

        # ── Hash cache: same file → instant ──────────────────────────────────
        if (st.session_state.get("_doc_hash") == file_hash
                and st.session_state.get("doc_name") == f.name
                and st.session_state.get("doc_ready")):
            step(100, "Document unchanged — ready instantly!")
            prog.empty()
            status.empty()
            st.session_state.file_path = tmp_path
            return True

        # Pre-warm embedding model NOW (non-blocking) so it's loaded by the
        # time the background FAISS thread starts.
        if not is_struct:
            from embeddings.embedding_model import get_embedding_model as _gw
            threading.Thread(target=_gw, daemon=True).start()

        # ── PHASE 1: Read + Chunk + BM25 + Prompts (blocking, fast) ──────────────
        step(10, "Reading document...")
        text = deps["load_document"](tmp_path)
        step(30, f"Loaded {len(text):,} characters")

        step(35, "Chunking text...")
        chunks = deps["chunk_text"](text)
        step(50, f"Created {len(chunks)} chunks")

        step(55, "Building search index...")
        # BM25 builds fast — wait for it
        deps["create_bm25_index"](chunks)
        step(85, "Index ready!")

        # Use default prompts immediately — generate real ones in background
        prompts = _default_prompts()
        def _gen_prompts_bg():
            try:
                result = gen_prompts(text)
                if result:
                    st.session_state.prompts = result
            except Exception:
                pass
        threading.Thread(target=_gen_prompts_bg, daemon=True).start()

        step(100, "Done — document ready!")
        st.session_state["_doc_hash"] = file_hash
        prog.empty()
        status.empty()

        # ── Unlock the UI immediately ──────────────────────────────────────────────────
        st.session_state.file_path  = tmp_path
        st.session_state.doc_name   = f.name
        st.session_state.doc_text   = text
        st.session_state.prompts    = prompts
        st.session_state.doc_ready  = True          # <── UI unlocks HERE
        st.session_state.chat       = []
        st.session_state.pending    = {}
        st.session_state.quiz_store = {}
        st.session_state.memory     = st.session_state.fmm.get_memory(f.name)

        existing = [d["name"] for d in st.session_state.multi_docs]
        if f.name not in existing:
            st.session_state.multi_docs.append(
                {"name": f.name, "path": tmp_path, "text": text, "chunks": chunks}
            )
        else:
            for d in st.session_state.multi_docs:
                if d["name"] == f.name:
                    d.update({"path": tmp_path, "text": text, "chunks": chunks})

        # ── PHASE 2: FAISS in background (non-blocking) ────────────────────────────────────────
        # UI is already live. FAISS embeds chunks silently; vector_search()
        # returns [] until it finishes, and hybrid_retriever uses BM25-only.
        if not is_struct:
            try:
                from vector_store.faiss_store import create_vector_store_background
                create_vector_store_background(chunks)
                logger.debug("FAISS background build launched")
            except Exception as e:
                logger.warning(f"FAISS background launch failed: {e}")

        return True

    except Exception as e:
        logger.error(f"process_doc: {e}")
        st.error(str(e))
        return False

def process_multi_parallel(files, deps):
    import tempfile
    from concurrent.futures import ThreadPoolExecutor, as_completed
    prog = st.progress(0); status = st.empty(); total = len(files); done = 0; all_docs = []
    def proc_one(f):
        suffix = "." + f.name.split(".")[-1]
        with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
            tmp.write(f.getvalue()); p = tmp.name
        ok, msg = deps["validate_file"](p)
        if not ok: return f.name, None, msg
        text = deps["load_document"](p); chunks = deps["chunk_text"](text)
        return f.name, {"name": f.name, "path": p, "text": text, "chunks": chunks}, None
    with ThreadPoolExecutor(max_workers=min(total, 3)) as ex:
        futs = {ex.submit(proc_one, f): f.name for f in files}
        for fut in as_completed(futs):
            nm, doc, err = fut.result(); done += 1
            prog.progress(int(done / total * 65))
            if err: status.markdown(f"<p style='color:#fca5a5;font-size:0.82rem;'>⚠ {nm}: {err}</p>", unsafe_allow_html=True)
            else:
                status.markdown(f"<p style='color:#a5b4fc;font-size:0.82rem;'>✓ {nm}</p>", unsafe_allow_html=True)
                all_docs.append(doc)
    if all_docs:
        prog.progress(70); status.markdown("<p style='color:#a5b4fc;font-size:0.82rem;'>⚙ Building indexes...</p>", unsafe_allow_html=True)
        active = all_docs[-1]; _build_indexes(active["chunks"], deps)
        existing = [d["name"] for d in st.session_state.multi_docs]
        for doc in all_docs:
            if doc["name"] not in existing: st.session_state.multi_docs.append(doc)
        prog.progress(90); status.markdown("<p style='color:#a5b4fc;font-size:0.82rem;'>⚙ Generating prompts...</p>", unsafe_allow_html=True)
        st.session_state.file_path = active["path"]; st.session_state.doc_name = active["name"]
        st.session_state.doc_text = active["text"]; st.session_state.prompts = gen_prompts(active["text"])
        st.session_state.doc_ready = True; st.session_state.chat = []
        st.session_state.pending = {}; st.session_state.quiz_store = {}
        st.session_state.memory = st.session_state.fmm.get_memory(active["name"])
    prog.progress(100); time.sleep(0.3); prog.empty(); status.empty()
    return len(all_docs)

def switch_document(doc, deps):
    """Switch active document. Indexes are built once and cached — instant after first visit."""
    doc_name  = doc["name"]
    cache_key = f"_idx_ready_{doc_name}"
 
    if not st.session_state.get(cache_key):
        # First visit: build indexes once, then cache so future switches are instant
        with st.spinner(f"Switching to {doc_name}..."):
            chunks = doc.get("chunks")
            if not chunks:
                chunks = deps["chunk_text"](doc["text"])
                doc["chunks"] = chunks
            _build_indexes(chunks, deps)
        st.session_state[cache_key] = True
    # If already cached: skip _build_indexes entirely — instant switch, no spinner
 
    st.session_state.doc_name   = doc["name"]
    st.session_state.file_path  = doc["path"]
    st.session_state.doc_text   = doc["text"]
    st.session_state.pending    = {}
    st.session_state.quiz_store = {}
    st.session_state.memory     = st.session_state.fmm.get_memory(doc["name"])



def answer_q(question, deps):
    st.session_state.pending = {}

    ok, msg = deps["validate_question"](question)
    if not ok:
        st.warning(msg); return
    if st.session_state.get("stop_requested"):
        st.session_state.stop_requested = False
        st.session_state.is_processing  = False
        return

    intent = detect_intent(question)

    # Safety net: N records + structured file → always regen_doc
    _fp  = st.session_state.get("file_path", "")
    _ext = os.path.splitext(_fp)[1].lower() if _fp else ""
    _q_toks = set(re.findall(r"[a-z0-9]+", question.lower()))
    _action_toks = {"generate", "create", "make", "give", "get", "export",
                    "produce", "download", "build", "prepare"}
    _file_toks = {"document", "doc", "file", "download", "downloadable"}
    _complaint_w = [
        "you provided", "you gave", "you only", "you said",
        "you told", "that's wrong", "not right", "incorrect",
        "not what i", "you made", "you only showed",
        "i told you", "you only generated", "you generated only",
        "only generated", "you missed", "you forgot",
    ]
    _is_complaint = any(w in question.lower() for w in _complaint_w)

    if (not intent["regen_doc"]
            and not _is_complaint
            and _ext in {".csv", ".xlsx", ".xls"}
            and re.search(r"\b\d+\s+(?:records?|rows?|entries|people|respondents?)\b", question.lower())
            and bool(_q_toks & _action_toks)
            and bool(_q_toks & _file_toks)):
        intent["regen_doc"] = True

    st.session_state.chat.append({"role": "human", "content": question})
    st.session_state.is_processing  = True
    st.session_state.stop_requested = False

    # ── Read conversation history ONCE — used by all paths ───────────────────
    _memory       = st.session_state.get("memory")
    _chat_history = ""
    try:
        if _memory and not _memory.is_empty():
            _chat_history = _memory.get_history_as_text()
    except Exception:
        _chat_history = ""

    # ── Quiz ──────────────────────────────────────────────────────────────────
    if intent["quiz"]:
        with st.spinner("Generating quiz..."):
            if st.session_state.get("stop_requested"):
                st.session_state.is_processing = False; return
            quiz = gen_quiz(st.session_state.doc_text)
        st.session_state.is_processing = False
        ks = uid(); st.session_state.quiz_store[ks] = quiz
        marker = f"__QUIZ__{ks}"
        st.session_state.chat.append({"role": "assistant", "content": marker})
        st.session_state.pending = {
            "evidence": [], "intent": intent,
            "question": question, "answer": marker, "ks": ks,
        }
        return

    # ── Document generation ───────────────────────────────────────────────────
    if intent["regen_doc"]:
        with st.spinner("Generating document..."):
            if st.session_state.get("stop_requested"):
                st.session_state.is_processing = False
                st.session_state.chat.pop(); return
            files_dict, generated_content = build_generated_document(
                st.session_state.doc_text, question, st.session_state.file_path,
            )
        st.session_state.is_processing = False
        _lm  = re.search(r'\b(\d+)\s+(?:records?|rows?)\b', question.lower())
        _n   = f"first {_lm.group(1)} records " if _lm else ""
        answer = f"✅ Document ready — {_n}use the download buttons below."
        st.session_state.chat.append({"role": "assistant", "content": answer})
        st.session_state.last_q = question
        st.session_state.last_a = answer
        st.session_state.pending = {
            "evidence": [], "intent": intent, "question": question,
            "answer": answer, "ks": uid(), "generated_files": files_dict,
        }
        return

    # ── All other queries ─────────────────────────────────────────────────────
    answer   = ""
    evidence = []
    _used_fast_path = False

    with st.spinner("Analysing..."):
        if st.session_state.get("stop_requested"):
            st.session_state.is_processing = False
            st.session_state.chat.pop(); return

        file_path     = st.session_state.get("file_path", "")
        file_ext      = os.path.splitext(file_path)[1].lower() if file_path else ""
        is_structured = file_ext in {".csv", ".xlsx", ".xls"}
        _special      = any(intent.get(k) for k in
                            ("excel", "chart", "pptx", "export_doc",
                             "regen_doc", "quiz", "resources"))

        _conceptual_words = [
            "explain", "what is", "what are", "describe", "tell me",
            "how does", "how is", "why ", "define", "meaning", "about",
            "summarize", "summary", "overview", "elaborate", "detail",
            "analyse", "analyze", "interpret", "simply", "simple",
            "for a", "like a", "like you", "child", "beginner",
            "understand", "help me", "teach", "educate",
            "what do you", "what can you", "who are you", "your role",
            "your purpose", "what are you", "generally", "introduce",
            "what will you", "how can you help", "capabilities",
        ]
        _has_possessive = bool(re.search(r"[A-Za-z]+'s\s+\w+", question))
        _conceptual     = (
            any(w in question.lower() for w in _conceptual_words)
            and not _has_possessive
        )

        if is_structured and file_path and not _special:
            if _conceptual:
                # Conceptual question → LLM with full context + history
                answer          = direct_summarize(
                    st.session_state.doc_text, question,
                    st.session_state.answer_mode, _chat_history,
                )
                _used_fast_path = True
            else:
                # Data query → dataframe agent
                try:
                    from agents.dataframe_agent import run_dataframe_agent
                    df_ans = run_dataframe_agent(question, file_path)
                    _bad   = {"", "none", "n/a", "could not compute"}
                    if (df_ans and df_ans.strip()
                            and df_ans.strip().lower() not in _bad
                            and not df_ans.lower().startswith("could not load")):
                        answer          = df_ans
                        _used_fast_path = True
                except Exception as df_err:
                    logger.warning(f"CSV fast-path: {df_err}")

        # ── Workflow for text docs or when fast-path returned nothing ─────────
        if not answer:
            try:
                result = deps["run_workflow"](
                    question=question,
                    memory=_memory,                       # workflow reads + writes memory
                    file_path=file_path,
                    answer_mode=st.session_state.answer_mode,
                    doc_text=st.session_state.doc_text,
                )
                if isinstance(result, dict):
                    answer   = result.get("answer", "") or result.get("output", "") or ""
                    evidence = result.get("evidence", []) or []
                else:
                    answer = str(result) if result else ""
                answer = answer.strip()
            except Exception as e:
                logger.error(f"workflow: {e}")

        # ── Fallback ──────────────────────────────────────────────────────────
        _empty = {"none", "result: none", "none.", "", "n/a"}
        if not answer or answer.strip().lower() in _empty or answer.startswith("Result: None"):
            answer = direct_summarize(
                st.session_state.doc_text, question,
                st.session_state.answer_mode, _chat_history,
            )
            _used_fast_path = True

        # ── Update memory for fast-path answers (workflow does its own) ───────
        if _used_fast_path and answer and answer.strip() and _memory:
            try:
                _memory.add_human_message(question)
                _memory.add_ai_message(answer)
            except Exception:
                pass

    if not answer or not answer.strip():
        answer = direct_summarize(
            st.session_state.doc_text, question,
            st.session_state.answer_mode, _chat_history,
        )

    # Clean up LLM artifacts
    import re as _re
    answer = _re.sub(r'^(Answer:|ANSWER:)\s*', '', answer.strip())
    answer = _re.sub(r'\n?\*?\*?Key Takeaway:?\*?\*?.*$', '', answer, flags=_re.IGNORECASE|_re.DOTALL).strip()

    st.session_state.is_processing = False
    st.session_state.chat.append({"role": "assistant", "content": answer})
    st.session_state.last_q = question
    st.session_state.last_a = answer
    st.session_state.pending = {
        "evidence": evidence, "intent": intent,
        "question": question, "answer": answer, "ks": uid(),
    }

def run_quick_action(action, deps):
    st.session_state.chat.append({"role": "human", "content": action})
    if "quiz" in action.lower():
        with st.spinner("Generating quiz..."):
            quiz = gen_quiz(st.session_state.doc_text)
        ks = uid(); st.session_state.quiz_store[ks] = quiz; marker = f"__QUIZ__{ks}"
        st.session_state.chat.append({"role": "assistant", "content": marker})
        st.session_state.pending = {"evidence": [], "intent": {"quiz": True}, "question": action, "answer": marker, "ks": ks}
        return
    with st.spinner("Working on it..."):
        result = direct_summarize(st.session_state.doc_text, action, st.session_state.answer_mode)
    if not result or not result.strip(): result = "Could not generate response. Please try again."
    st.session_state.chat.append({"role": "assistant", "content": result})
    st.session_state.last_q = action; st.session_state.last_a = result
    st.session_state.pending = {"evidence": [], "intent": {}, "question": action, "answer": result, "ks": uid()}

def render_pending(p):
    if not p:
        return

    ev       = p.get("evidence", [])
    intent   = p.get("intent", {})
    question = p.get("question", "")
    answer   = p.get("answer", "")
    ks       = p.get("ks", "")

    # Evidence
    if ev:
        with st.expander("Evidence Sources", expanded=False):
            for i, chunk in enumerate(ev, 1):
                st.markdown(
                    f"<div style='background:rgba(99,102,241,0.04);border-left:2px solid "
                    f"rgba(99,102,241,0.3);border-radius:0 8px 8px 0;padding:10px 14px;"
                    f"margin:6px 0;font-size:0.79rem;color:#64748b;'>"
                    f"<strong style='color:#a5b4fc;'>Source {i}</strong><br>{chunk}</div>",
                    unsafe_allow_html=True,
                )

    if intent.get("quiz"):
        return

    # ── Charts (cached) ───────────────────────────────────────────────────────
    if intent.get("chart"):
        _k = f"__chart_{ks}"
        if _k not in st.session_state:
            with st.spinner("Generating chart..."):
                st.session_state[_k] = gen_charts(st.session_state.doc_text, question)
        charts = st.session_state[_k]
        if charts:
            c1, c2 = st.columns(2) if len(charts) > 1 else (st.container(), None)
            for i, (_, fig) in enumerate(charts[:4]):
                with (c1 if i % 2 == 0 else c2) if c2 else c1:
                    st.plotly_chart(fig, width="stretch")
        else:
            st.info("No numerical data found for charting. Specify a column name.")
        return

    # ── Excel (cached) ────────────────────────────────────────────────────────
    if intent.get("excel"):
        st.markdown(
            "<span style='font-size:0.62rem;font-weight:700;color:rgba(16,185,129,0.8);"
            "text-transform:uppercase;letter-spacing:2px;'>Excel Export</span>",
            unsafe_allow_html=True,
        )
        _k = f"__excel_{ks}"
        if _k not in st.session_state:
            with st.spinner("Creating Excel..."):
                st.session_state[_k] = gen_excel_from_doc(st.session_state.doc_text, question)
        st.download_button("Download Excel", data=st.session_state[_k],
                           file_name="documind_data.xlsx",
                           mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                           key=f"dlx_{ks}")
        return

    # ── PowerPoint (cached) ───────────────────────────────────────────────────
    if intent.get("pptx"):
        _k = f"__pptx_{ks}"
        if _k not in st.session_state:
            with st.spinner("Creating PowerPoint..."):
                st.session_state[_k] = gen_pptx(answer, st.session_state.doc_name)
        pb = st.session_state[_k]
        if pb:
            st.download_button("Download PowerPoint", data=pb, file_name="documind.pptx",
                               mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                               key=f"dlp_{ks}")
        return

    # ── Regenerated document (cached) ─────────────────────────────────────────
    if intent.get("regen_doc"):
        _k = f"__regen_{ks}"
        files_dict = p.get("generated_files")
        if not files_dict:
            if _k not in st.session_state:
                with st.spinner("Generating document..."):
                    st.session_state[_k], _ = build_generated_document(
                        st.session_state.doc_text, question, st.session_state.file_path
                    )
            files_dict = st.session_state[_k]

        st.markdown(
            "<span style='font-size:0.62rem;font-weight:700;color:rgba(99,102,241,0.8);"
            "text-transform:uppercase;letter-spacing:2px;'>Regenerated Document</span>",
            unsafe_allow_html=True,
        )
        c1, c2, c3, c4 = st.columns(4)
        with c1: st.download_button("DOCX", data=files_dict["docx"],
                                    file_name="documind_regen.docx", key=f"rg_docx_{ks}",
                                    mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document")
        with c2: st.download_button("PDF",  data=files_dict["pdf"],
                                    file_name="documind_regen.pdf",  key=f"rg_pdf_{ks}",
                                    mime="application/pdf")
        with c3: st.download_button("CSV",  data=files_dict["csv"],
                                    file_name="documind_regen.csv",  key=f"rg_csv_{ks}",
                                    mime="text/csv")
        with c4: st.download_button("TXT",  data=files_dict["txt"],
                                    file_name="documind_regen.txt",  key=f"rg_txt_{ks}",
                                    mime="text/plain")
        return

    # ── Export answer as file ──────────────────────────────────────────────────
    if intent.get("export_doc"):
        st.markdown(
            "<span style='font-size:0.62rem;font-weight:700;color:rgba(16,185,129,0.8);"
            "text-transform:uppercase;letter-spacing:2px;'>Download Answer</span>",
            unsafe_allow_html=True,
        )
        cols = st.columns(5)
        for i, (fmt, lbl, mime) in enumerate([
            ("txt",  "TXT",   "text/plain"),
            ("docx", "DOCX",  "application/vnd.openxmlformats-officedocument.wordprocessingml.document"),
            ("pdf",  "PDF",   "application/pdf"),
            ("csv",  "CSV",   "text/csv"),
            ("xlsx", "Excel", "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"),
        ]):
            with cols[i]:
                st.download_button(lbl, data=gen_file(answer, fmt),
                                   file_name=f"documind.{fmt}", mime=mime,
                                   key=f"dl_{fmt}_{ks}")

    # ── Resources ──────────────────────────────────────────────────────────────
    if intent.get("resources"):
        show_resources(question, answer)



CSS = """<style>
html,body,[class*="css"]{font-family:-apple-system,BlinkMacSystemFont,'Segoe UI','Helvetica Neue',Arial,sans-serif !important;}
.stApp{background:radial-gradient(ellipse at 20% 20%,rgba(99,102,241,0.05) 0%,transparent 55%),radial-gradient(ellipse at 80% 80%,rgba(6,182,212,0.04) 0%,transparent 55%),#04060f !important;}
section[data-testid="stSidebar"]{background:rgba(4,6,15,0.98) !important;border-right:1px solid rgba(99,102,241,0.1) !important;}
.main .block-container{padding-left:1rem !important;padding-right:1rem !important;padding-top:0.5rem !important;max-width:100% !important;}
section[data-testid="stSidebar"]+.main .block-container{padding-left:1rem !important;padding-right:1rem !important;}
div[data-testid="stChatMessageContainer"]{max-width:100% !important;}
div[data-testid="stChatMessage"]{max-width:100% !important;}
section[data-testid="stSidebar"] *{color:#94a3b8 !important;}
.stButton>button{background:linear-gradient(135deg,#4f46e5,#6366f1) !important;color:#fff !important;border:1px solid rgba(99,102,241,0.4) !important;border-radius:10px !important;font-weight:600 !important;font-size:0.79rem !important;transition:all 0.2s !important;box-shadow:0 2px 10px rgba(79,70,229,0.2) !important;}
.stButton>button:hover{background:linear-gradient(135deg,#6366f1,#818cf8) !important;transform:translateY(-2px) !important;box-shadow:0 5px 18px rgba(99,102,241,0.35) !important;}
.stChatInputContainer{background:rgba(8,10,22,0.9) !important;border:1px solid rgba(99,102,241,0.2) !important;border-radius:16px !important;box-shadow:0 4px 20px rgba(0,0,0,0.25) !important;}
.stChatInputContainer:focus-within{border-color:rgba(99,102,241,0.45) !important;box-shadow:0 0 0 3px rgba(99,102,241,0.09) !important;}
.stChatInputContainer textarea{background:transparent !important;color:#e2e8f0 !important;}
div[data-testid="stChatMessage"]{background:transparent !important;border-radius:14px !important;}
div[data-testid="stExpander"]{background:rgba(8,10,22,0.6) !important;border:1px solid rgba(99,102,241,0.1) !important;border-radius:12px !important;}
div[data-testid="stSelectbox"]>div>div{background:rgba(8,10,22,0.9) !important;border:1px solid rgba(99,102,241,0.18) !important;border-radius:10px !important;}
section[data-testid="stFileUploaderDropzone"]{background:rgba(99,102,241,0.02) !important;border:2px dashed rgba(99,102,241,0.25) !important;border-radius:14px !important;}
div[data-testid="stProgressBar"]>div>div{background:linear-gradient(90deg,#4f46e5,#7c3aed,#06b6d4) !important;border-radius:6px !important;box-shadow:0 0 10px rgba(99,102,241,0.35) !important;}
div[data-testid="stSuccess"]{background:rgba(16,185,129,0.07) !important;border:1px solid rgba(16,185,129,0.22) !important;border-radius:11px !important;}
div[data-testid="stError"]{background:rgba(239,68,68,0.07) !important;border:1px solid rgba(239,68,68,0.22) !important;border-radius:11px !important;}
div[data-testid="stInfo"]{background:rgba(99,102,241,0.06) !important;border:1px solid rgba(99,102,241,0.18) !important;border-radius:11px !important;}
div[data-testid="stWarning"]{background:rgba(245,158,11,0.06) !important;border:1px solid rgba(245,158,11,0.2) !important;border-radius:11px !important;}
div[data-testid="stDownloadButton"]>button{background:rgba(16,185,129,0.09) !important;border:1px solid rgba(16,185,129,0.25) !important;color:#6ee7b7 !important;border-radius:10px !important;font-weight:600 !important;}
div[data-testid="stDownloadButton"]>button:hover{background:rgba(16,185,129,0.16) !important;transform:translateY(-1px) !important;}
div[data-testid="stRadio"] label{background:rgba(8,10,22,0.55) !important;border:1px solid rgba(255,255,255,0.06) !important;border-radius:9px !important;padding:8px 14px !important;margin:3px 0 !important;color:#94a3b8 !important;}
::-webkit-scrollbar{width:3px;height:3px;}::-webkit-scrollbar-track{background:transparent;}::-webkit-scrollbar-thumb{background:rgba(99,102,241,0.22);border-radius:4px;}
hr{border-color:rgba(255,255,255,0.04) !important;margin:10px 0 !important;}
/* Answer alignment: scale down large headings inside chat messages */
div[data-testid="stChatMessage"] h1{font-size:1.25rem !important;margin:12px 0 6px !important;color:#e2e8f0 !important;}
div[data-testid="stChatMessage"] h2{font-size:1.1rem !important;margin:10px 0 5px !important;color:#e2e8f0 !important;}
div[data-testid="stChatMessage"] h3{font-size:1rem !important;margin:8px 0 4px !important;color:#e2e8f0 !important;}
div[data-testid="stChatMessage"] h4{font-size:0.9rem !important;margin:6px 0 3px !important;color:#e2e8f0 !important;}
div[data-testid="stChatMessage"] p{margin:4px 0 !important;line-height:1.65 !important;}
div[data-testid="stChatMessage"] ul,div[data-testid="stChatMessage"] ol{margin:4px 0 4px 18px !important;padding-left:0 !important;}
div[data-testid="stChatMessage"] li{margin:3px 0 !important;line-height:1.55 !important;}
div[data-testid="stChatMessage"] strong{color:#c7d2fe !important;}
.stop-btn>button{background:rgba(239,68,68,0.1) !important;border:1px solid rgba(239,68,68,0.25) !important;color:#fca5a5 !important;border-radius:9px !important;font-size:0.78rem !important;box-shadow:none !important;}
.stop-btn>button:hover{background:rgba(239,68,68,0.18) !important;border-color:rgba(239,68,68,0.4) !important;transform:none !important;}
</style>"""

SIDEBAR_CSS = """<style>
.dm-logo{display:flex;align-items:center;gap:12px;padding:18px 0 16px;border-bottom:1px solid rgba(99,102,241,0.1);margin-bottom:18px;}
/* .dm-icon replaced by SVG logo */
.dm-name{font-size:1rem !important;font-weight:800 !important;color:#f1f5f9 !important;line-height:1.2;}
.dm-sub{font-size:0.56rem !important;color:rgba(99,102,241,0.7) !important;font-weight:700 !important;text-transform:uppercase;letter-spacing:2.5px;}
.sec-lbl{font-size:0.56rem !important;font-weight:800 !important;color:rgba(99,102,241,0.5) !important;text-transform:uppercase !important;letter-spacing:3px !important;display:block !important;margin:14px 0 8px !important;}
.doc-card{background:linear-gradient(135deg,rgba(99,102,241,0.06),rgba(6,182,212,0.03));border:1px solid rgba(99,102,241,0.13);border-radius:11px;padding:9px 13px;margin:5px 0;}
.doc-name{color:#a5b4fc !important;font-size:0.78rem !important;font-weight:600 !important;white-space:nowrap;overflow:hidden;text-overflow:ellipsis;max-width:195px;}
.doc-size{color:#3f4f6b !important;font-size:0.66rem !important;margin-top:2px;}
.msg-c{background:linear-gradient(135deg,rgba(99,102,241,0.07),rgba(6,182,212,0.04));border:1px solid rgba(99,102,241,0.11);border-radius:11px;padding:13px;text-align:center;margin:6px 0;}
.mc-n{font-size:2rem !important;font-weight:900 !important;background:linear-gradient(135deg,#a5b4fc,#06b6d4) !important;-webkit-background-clip:text !important;-webkit-text-fill-color:transparent !important;background-clip:text !important;display:block !important;line-height:1 !important;}
.mc-l{font-size:0.56rem !important;color:#3f4f6b !important;text-transform:uppercase !important;letter-spacing:2.5px !important;margin-top:4px !important;display:block !important;}
</style>"""




@st.cache_data(show_spinner=False)
def _hero_html():
    return """<!DOCTYPE html><html><head><meta charset="UTF-8"><style>
*{margin:0;padding:0;box-sizing:border-box;font-family:-apple-system,BlinkMacSystemFont,Segoe UI,sans-serif;}
body{background:radial-gradient(ellipse at 30% 20%,rgba(99,102,241,0.07) 0%,transparent 55%),radial-gradient(ellipse at 70% 80%,rgba(6,182,212,0.05) 0%,transparent 55%),#04060f;min-height:100vh;display:flex;flex-direction:column;align-items:center;justify-content:center;padding:48px 24px;text-align:center;}
@keyframes blink{0%,100%{opacity:1}50%{opacity:0.15}}
@keyframes flow{0%{background-position:0% center}100%{background-position:300% center}}
@keyframes up{from{opacity:0;transform:translateY(18px)}to{opacity:1;transform:translateY(0)}}
.badge{display:inline-flex;align-items:center;gap:8px;background:rgba(99,102,241,0.08);border:1px solid rgba(99,102,241,0.22);border-radius:100px;padding:7px 20px;font-size:0.61rem;font-weight:700;color:#818cf8;letter-spacing:3px;text-transform:uppercase;margin-bottom:30px;animation:up 0.4s ease forwards;}
.dot{width:5px;height:5px;background:#22c55e;border-radius:50%;animation:blink 2s ease infinite;box-shadow:0 0 6px #22c55e;}
h1{font-size:clamp(2.6rem,7vw,5.4rem);font-weight:900;letter-spacing:-3.5px;line-height:0.96;margin-bottom:20px;background:linear-gradient(135deg,#fff 0%,#c7d2fe 20%,#a5b4fc 40%,#06b6d4 65%,#a5b4fc 80%,#fff 100%);background-size:300% auto;-webkit-background-clip:text;-webkit-text-fill-color:transparent;background-clip:text;animation:flow 5s linear infinite,up 0.6s ease forwards;}
.sub{font-size:0.98rem;color:#374151;max-width:440px;line-height:1.85;margin:0 auto 44px;animation:up 0.8s ease forwards;}
.sub b{color:#6366f1;font-weight:600;}
.stats{display:flex;border:1px solid rgba(255,255,255,0.06);border-radius:14px;overflow:hidden;background:rgba(8,10,22,0.85);max-width:520px;margin:0 auto 52px;animation:up 1s ease forwards;}
.s{flex:1;padding:18px 8px;text-align:center;border-right:1px solid rgba(255,255,255,0.05);}
.s:last-child{border-right:none;}
.sv{font-size:1.5rem;font-weight:900;background:linear-gradient(135deg,#a5b4fc,#06b6d4);-webkit-background-clip:text;-webkit-text-fill-color:transparent;background-clip:text;display:block;}
.sl{font-size:0.55rem;color:#1e293b;text-transform:uppercase;letter-spacing:2px;margin-top:3px;display:block;}
.grid{display:grid;grid-template-columns:repeat(4,1fr);gap:10px;max-width:880px;margin:0 auto;animation:up 1.2s ease forwards;}
.card{background:rgba(8,10,22,0.7);border:1px solid rgba(255,255,255,0.05);border-radius:13px;padding:18px 14px;text-align:left;transition:all 0.28s;}
.card:hover{border-color:rgba(99,102,241,0.28);transform:translateY(-4px);box-shadow:0 14px 35px rgba(99,102,241,0.13);}
.num{font-size:0.57rem;font-weight:800;color:rgba(99,102,241,0.38);letter-spacing:3px;text-transform:uppercase;margin-bottom:8px;}
.ico{font-size:1rem;margin-bottom:7px;display:block;}
.nm{color:#e2e8f0;font-size:0.82rem;font-weight:700;margin-bottom:4px;}
.tx{color:#1e293b;font-size:0.7rem;line-height:1.55;}
</style></head><body>
<div class="badge"><div class="dot"></div>Enterprise Document Intelligence</div>
<h1>DocuMind AI</h1>
<p class="sub">Upload any document. Ask in <b>plain English</b>.<br>Get precise answers, charts, quizzes, and exports instantly.</p>
<div class="stats">
<div class="s"><span class="sv">6+</span><span class="sl">Formats</span></div>
<div class="s"><span class="sv">5</span><span class="sl">AI Agents</span></div>
<div class="s"><span class="sv">20+</span><span class="sl">Features</span></div>
<div class="s"><span class="sv">100%</span><span class="sl">Private</span></div></div>
<div class="grid">
<div class="card"><div class="num">01</div><span class="ico">📊</span><div class="nm">Auto Charts</div><div class="tx">Ask "show chart" for interactive Plotly graphs</div></div>
<div class="card"><div class="num">02</div><span class="ico">📝</span><div class="nm">Quiz Mode</div><div class="tx">Ask "give me a quiz" for MCQ with instant checking</div></div>
<div class="card"><div class="num">03</div><span class="ico">🔊</span><div class="nm">Read Aloud</div><div class="tx">Browser TTS reads every answer aloud</div></div>
<div class="card"><div class="num">04</div><span class="ico">🔗</span><div class="nm">Live Resources</div><div class="tx">Ask "show resources" for YouTube, Google, Scholar</div></div>
<div class="card"><div class="num">05</div><span class="ico">⬇</span><div class="nm">Export Any Format</div><div class="tx">Ask "export as pdf/docx/excel" to download</div></div>
<div class="card"><div class="num">06</div><span class="ico">🔄</span><div class="nm">Regenerate Docs</div><div class="tx">Ask "shorten to 10 records and regenerate"</div></div>
<div class="card"><div class="num">07</div><span class="ico">🔍</span><div class="nm">Hybrid Search</div><div class="tx">FAISS semantic + BM25 keyword fusion</div></div>
<div class="card"><div class="num">08</div><span class="ico">🗂</span><div class="nm">Multi-Document</div><div class="tx">Upload multiple files, switch, compare all</div></div>
</div></body></html>"""

def show_hero():
    st.iframe(_hero_html(), height=880)




def query_multi_direct(question, docs):
    """Query all documents and synthesize a unified answer."""
    combined = ""
    for d in docs[:5]:
        combined += f"\n\n{'='*50}\nDOCUMENT: {d['name']}\n{'='*50}\n{d['text'][:2500]}"
    tpl = """You are DocuMind AI. Answer this question using ALL the documents below.

{docs}

Question: {question}

RULES:
- Reference specific documents by name when citing facts
- Compare/contrast information across documents where relevant
- If information appears in multiple documents, say so
- Be comprehensive and specific
- End with: **Summary:** [one sentence covering all documents]

Answer:"""
    try: return llm_call(tpl, {"docs": combined, "question": question})
    except Exception as e: return f"Error querying documents: {e}"

def compare_docs_direct(docs):
    """Generate detailed comparison across all documents."""
    combined = ""
    for d in docs[:5]:
        combined += f"\n\n{'='*50}\nDOCUMENT: {d['name']}\n{'='*50}\n{d['text'][:1800]}"
    tpl = """Compare ALL these documents comprehensively:

{docs}

Create a detailed comparison covering:
## 1. Document Summaries
[Brief summary of each document]

## 2. Key Similarities
[What they have in common]

## 3. Key Differences
[How they differ]

## 4. Unique Information Per Document
[What each document uniquely contains]

## 5. Overall Conclusion
[What insights emerge from reading all together]

Use markdown formatting with tables where helpful."""
    try: return llm_call(tpl, {"docs": combined})
    except Exception as e: return f"Comparison error: {e}"




def _show_upload_sidebar_only():
    """Lightweight sidebar shown on hero page — zero heavy imports."""
    with st.sidebar:
        render_startup_warnings()
        st.markdown("""
<div class='dm-logo'>
  <div style='flex-shrink:0;'>
    <svg width="46" height="46" viewBox="0 0 46 46" fill="none" xmlns="http://www.w3.org/2000/svg">
      <defs>
        <linearGradient id="bgGrad" x1="0" y1="0" x2="46" y2="46" gradientUnits="userSpaceOnUse">
          <stop offset="0%" stop-color="#0d0b1e"/>
          <stop offset="30%" stop-color="#2d2a8a"/>
          <stop offset="65%" stop-color="#5b21b6"/>
          <stop offset="100%" stop-color="#0e7490"/>
        </linearGradient>
        <linearGradient id="shineGrad" x1="0" y1="0" x2="0" y2="1">
          <stop offset="0%" stop-color="white" stop-opacity="0.13"/>
          <stop offset="100%" stop-color="white" stop-opacity="0"/>
        </linearGradient>
        <radialGradient id="nodeGrad" cx="50%" cy="50%" r="50%">
          <stop offset="0%" stop-color="#22d3ee"/>
          <stop offset="100%" stop-color="#0891b2"/>
        </radialGradient>
        <filter id="nodeGlow" x="-80%" y="-80%" width="260%" height="260%">
          <feGaussianBlur in="SourceGraphic" stdDeviation="2.5" result="blur"/>
          <feComposite in="SourceGraphic" in2="blur" operator="over"/>
        </filter>
        <filter id="softGlow" x="-40%" y="-40%" width="180%" height="180%">
          <feGaussianBlur in="SourceGraphic" stdDeviation="1.2" result="blur"/>
          <feMerge><feMergeNode in="blur"/><feMergeNode in="SourceGraphic"/></feMerge>
        </filter>
        <clipPath id="rounded"><rect width="46" height="46" rx="12"/></clipPath>
      </defs>
      <rect x="1" y="1" width="44" height="44" rx="12" fill="#4f46e5" opacity="0.35" filter="url(#nodeGlow)"/>
      <rect width="46" height="46" rx="12" fill="url(#bgGrad)"/>
      <rect width="46" height="26" rx="12" fill="url(#shineGrad)" clip-path="url(#rounded)"/>
      <rect x="0.75" y="0.75" width="44.5" height="44.5" rx="11.5" fill="none" stroke="white" stroke-opacity="0.1" stroke-width="1.5"/>
      <rect x="10" y="11" width="3.5" height="24" rx="1.75" fill="white" opacity="0.92"/>
      <path d="M 13.5 11.5 C 26 11.5 32.5 16.5 32.5 23 C 32.5 29.5 26 34.5 13.5 34.5" fill="none" stroke="white" stroke-width="3" stroke-linecap="round" opacity="0.88"/>
      <rect x="17" y="17.5" width="8" height="1.8" rx="0.9" fill="url(#nodeGrad)" opacity="0.9"/>
      <rect x="17" y="22" width="11" height="1.8" rx="0.9" fill="url(#nodeGrad)" opacity="0.72"/>
      <rect x="17" y="26.5" width="7" height="1.8" rx="0.9" fill="url(#nodeGrad)" opacity="0.55"/>
      <circle cx="37" cy="12" r="5.5" fill="#06b6d4" opacity="0.12" filter="url(#nodeGlow)"/>
      <circle cx="37" cy="12" r="4" fill="none" stroke="#22d3ee" stroke-width="1.2" stroke-opacity="0.6"/>
      <circle cx="37" cy="12" r="2.6" fill="url(#nodeGrad)" filter="url(#softGlow)" opacity="0.98"/>
      <circle cx="37" cy="12" r="1.2" fill="white" opacity="0.9"/>
      <line x1="31" y1="15" x2="34" y2="13.5" stroke="#22d3ee" stroke-width="1" stroke-opacity="0.4" stroke-dasharray="1.5 1.5"/>
      <rect x="10" y="37.5" width="26" height="1.2" rx="0.6" fill="url(#nodeGrad)" opacity="0.35"/>
    </svg>
  </div>
  <div><div class='dm-name'>DocuMind AI</div><div class='dm-sub'>Document Intelligence</div></div>
</div>""", unsafe_allow_html=True)
        st.markdown("<span class='sec-lbl'>Upload Document(s)</span>", unsafe_allow_html=True)
        files = st.file_uploader(
            "upload", type=["pdf","docx","pptx","ppt","xlsx","csv","txt"],
            accept_multiple_files=True, label_visibility="collapsed",
        )
        if files:
            for f in files:
                st.markdown(
                    f"<div class='doc-card'><div class='doc-name'>{f.name}</div>"
                    f"<div class='doc-size'>{f.size/1024:.0f} KB</div></div>",
                    unsafe_allow_html=True)
            lbl = f"Process {len(files)} Files" if len(files) > 1 else "Process Document"
            if st.button(lbl, width="stretch"):

                st.session_state["_process_clicked"] = True
                st.session_state["_pending_files"] = files
                st.rerun()

def show_startup_check():
    if st.session_state.get("startup_checked"):
        return
    try:
        from utils.config import validate_startup
        results = validate_startup()
        st.session_state["startup_results"] = results
        st.session_state["startup_checked"] = True
    except Exception as e:
        logger.error(f"Startup check failed: {e}")
        st.session_state["startup_checked"] = True
        st.session_state["startup_results"] = {"errors": [str(e)], "warnings": []}


def render_startup_warnings():
    results = st.session_state.get("startup_results", {})
    for err in results.get("errors", []):
        st.sidebar.error(f"Setup Error: {err}")
    for warn in results.get("warnings", []):
        st.sidebar.warning(warn)


def main():

    st.markdown(CSS, unsafe_allow_html=True)
    st.markdown(SIDEBAR_CSS, unsafe_allow_html=True)




    _doc_ready = st.session_state.get("doc_ready", False)
    _process_clicked = st.session_state.get("_process_clicked", False)

    if not _doc_ready and not _process_clicked:

        _show_upload_sidebar_only()
        show_hero()
        return


    deps = load_core()
    show_startup_check()
    init(deps)


    if st.session_state.get("_process_clicked") and not st.session_state.get("doc_ready"):
        pending = st.session_state.get("_pending_files", [])
        st.session_state["_process_clicked"] = False
        st.session_state["_pending_files"] = []
        if pending:
            if len(pending) == 1:
                if process_doc(pending[0], deps):
                    st.success(f"Ready: {pending[0].name}")
                    st.rerun()
            else:
                n = process_multi_parallel(pending, deps)
                if n > 0:
                    st.success(f"{n} documents ready")
                    st.rerun()

    with st.sidebar:
        st.markdown("""
<div class='dm-logo'>
  <div style='flex-shrink:0;'>
    <svg width="46" height="46" viewBox="0 0 46 46" fill="none" xmlns="http://www.w3.org/2000/svg">
      <defs>
        <linearGradient id="bgGrad" x1="0" y1="0" x2="46" y2="46" gradientUnits="userSpaceOnUse">
          <stop offset="0%" stop-color="#0d0b1e"/>
          <stop offset="30%" stop-color="#2d2a8a"/>
          <stop offset="65%" stop-color="#5b21b6"/>
          <stop offset="100%" stop-color="#0e7490"/>
        </linearGradient>
        <linearGradient id="shineGrad" x1="0" y1="0" x2="0" y2="1">
          <stop offset="0%" stop-color="white" stop-opacity="0.13"/>
          <stop offset="100%" stop-color="white" stop-opacity="0"/>
        </linearGradient>
        <radialGradient id="nodeGrad" cx="50%" cy="50%" r="50%">
          <stop offset="0%" stop-color="#22d3ee"/>
          <stop offset="100%" stop-color="#0891b2"/>
        </radialGradient>
        <filter id="nodeGlow" x="-80%" y="-80%" width="260%" height="260%">
          <feGaussianBlur in="SourceGraphic" stdDeviation="2.5" result="blur"/>
          <feComposite in="SourceGraphic" in2="blur" operator="over"/>
        </filter>
        <filter id="softGlow" x="-40%" y="-40%" width="180%" height="180%">
          <feGaussianBlur in="SourceGraphic" stdDeviation="1.2" result="blur"/>
          <feMerge><feMergeNode in="blur"/><feMergeNode in="SourceGraphic"/></feMerge>
        </filter>
        <clipPath id="rounded"><rect width="46" height="46" rx="12"/></clipPath>
      </defs>
      <rect x="1" y="1" width="44" height="44" rx="12" fill="#4f46e5" opacity="0.35" filter="url(#nodeGlow)"/>
      <rect width="46" height="46" rx="12" fill="url(#bgGrad)"/>
      <rect width="46" height="26" rx="12" fill="url(#shineGrad)" clip-path="url(#rounded)"/>
      <rect x="0.75" y="0.75" width="44.5" height="44.5" rx="11.5" fill="none" stroke="white" stroke-opacity="0.1" stroke-width="1.5"/>
      <rect x="10" y="11" width="3.5" height="24" rx="1.75" fill="white" opacity="0.92"/>
      <path d="M 13.5 11.5 C 26 11.5 32.5 16.5 32.5 23 C 32.5 29.5 26 34.5 13.5 34.5" fill="none" stroke="white" stroke-width="3" stroke-linecap="round" opacity="0.88"/>
      <rect x="17" y="17.5" width="8" height="1.8" rx="0.9" fill="url(#nodeGrad)" opacity="0.9"/>
      <rect x="17" y="22" width="11" height="1.8" rx="0.9" fill="url(#nodeGrad)" opacity="0.72"/>
      <rect x="17" y="26.5" width="7" height="1.8" rx="0.9" fill="url(#nodeGrad)" opacity="0.55"/>
      <circle cx="37" cy="12" r="5.5" fill="#06b6d4" opacity="0.12" filter="url(#nodeGlow)"/>
      <circle cx="37" cy="12" r="4" fill="none" stroke="#22d3ee" stroke-width="1.2" stroke-opacity="0.6"/>
      <circle cx="37" cy="12" r="2.6" fill="url(#nodeGrad)" filter="url(#softGlow)" opacity="0.98"/>
      <circle cx="37" cy="12" r="1.2" fill="white" opacity="0.9"/>
      <line x1="31" y1="15" x2="34" y2="13.5" stroke="#22d3ee" stroke-width="1" stroke-opacity="0.4" stroke-dasharray="1.5 1.5"/>
      <rect x="10" y="37.5" width="26" height="1.2" rx="0.6" fill="url(#nodeGrad)" opacity="0.35"/>
    </svg>
  </div>
  <div><div class='dm-name'>DocuMind AI</div><div class='dm-sub'>Document Intelligence</div></div>
</div>""", unsafe_allow_html=True)
        st.markdown("<span class='sec-lbl'>Upload Document(s)</span>", unsafe_allow_html=True)
        files = st.file_uploader("upload", type=["pdf","docx","pptx","ppt","xlsx","csv","txt"],
                                  accept_multiple_files=True, label_visibility="collapsed")
        if files:
            for f in files:
                st.markdown(f"<div class='doc-card'><div class='doc-name'>{f.name}</div><div class='doc-size'>{f.size/1024:.0f} KB</div></div>", unsafe_allow_html=True)
            lbl = f"Process {len(files)} Files" if len(files) > 1 else "Process Document"
            if st.button(lbl, width="stretch"):
                if len(files) == 1:
                    if process_doc(files[0], deps): st.success(f"Ready: {files[0].name}"); st.rerun()
                else:
                    n = process_multi_parallel(files, deps)
                    if n > 0: st.success(f"{n} documents ready"); st.rerun()

        if st.session_state.doc_ready:
            st.markdown("---")
            st.markdown(
                "<span style='font-size:0.62rem;font-weight:700;color:rgba(99,102,241,0.6);"
                "text-transform:uppercase;letter-spacing:2px;'>Response Length</span>",
                unsafe_allow_html=True,
            )
            mode = st.selectbox(
                "Response Length",
                ["detailed", "quick", "bullet", "executive", "beginner"],
                index=["detailed", "quick", "bullet", "executive", "beginner"].index(
                    st.session_state.answer_mode if st.session_state.answer_mode in ["detailed", "quick", "bullet", "executive", "beginner"] else "detailed"
                ),
                label_visibility="collapsed",
                key="answer_mode",
            )
            st.markdown("---")
            st.markdown("<span class='sec-lbl'>Active Document</span>", unsafe_allow_html=True)
            st.markdown(f"<div class='doc-card'><div class='doc-name'>📄 {st.session_state.doc_name}</div></div>", unsafe_allow_html=True)
            st.markdown(f"<div class='msg-c'><span class='mc-n'>{len(st.session_state.chat)}</span><span class='mc-l'>Messages</span></div>", unsafe_allow_html=True)
            if len(st.session_state.multi_docs) > 1:
                st.markdown("---")
                st.markdown("<span class='sec-lbl'>All Documents</span>", unsafe_allow_html=True)
                for i, doc in enumerate(st.session_state.multi_docs):
                    is_active = doc["name"] == st.session_state.doc_name
                    color  = "#3fb950" if is_active else "#64748b"
                    marker = "●" if is_active else "○"
                    st.markdown(
                        f"<div style='color:{color};font-size:0.77rem;padding:3px 0;"
                        f"overflow:hidden;text-overflow:ellipsis;white-space:nowrap;'>"
                        f"{marker} {doc['name']}</div>",
                        unsafe_allow_html=True,
                    )
                    if not is_active:
                        if st.button("Use", key=f"sw_{i}", use_container_width=True):
                            switch_document(doc, deps)
                            st.rerun()
            st.markdown("---")
            st.markdown("<span class='sec-lbl'>Quick Actions</span>", unsafe_allow_html=True)
            for i, act in enumerate(["Summarize document","Extract action items","Identify all risks","Extract key metrics","Generate FAQ","Create a quiz"]):
                if st.button(act, width="stretch", key=f"qa_{i}"):
                    run_quick_action(act, deps); st.rerun()
            st.markdown("---")

            st.markdown("---")
            cc1, cc2 = st.columns(2)
            with cc1:
                if st.button("Clear", width="stretch"):
                    st.session_state.memory.clear(); st.session_state.chat = []
                    st.session_state.pending = {}; st.session_state.quiz_store = {}; st.rerun()
            with cc2:
                if st.button("New File", width="stretch"):
                    for k in ["doc_ready","file_path","doc_name","doc_text","chat","prompts","multi_docs","pending","quiz_store","is_processing","stop_requested"]:
                        st.session_state[k] = ([] if k in ["chat","prompts","multi_docs"] else {} if k in ["pending","quiz_store"] else False if k in ["doc_ready","is_processing","stop_requested"] else "")
                    st.session_state.memory = deps["SessionMemory"](); st.rerun()

    if not st.session_state.doc_ready:
        show_hero(); return

    if st.session_state.prompts:
        st.markdown("<p style='color:rgba(99,102,241,0.55);font-size:0.61rem;font-weight:800;text-transform:uppercase;letter-spacing:3px;margin-bottom:10px;'>Suggested</p>", unsafe_allow_html=True)
        cols = st.columns(3)
        for i, prompt in enumerate(st.session_state.prompts):
            with cols[i % 3]:
                if st.button(prompt, key=f"sp_{i}", width="stretch"):
                    answer_q(prompt, deps); st.rerun()
        st.markdown("---")

    if len(st.session_state.multi_docs) > 1:
        with st.expander(f"Multi-Document Panel  -  {len(st.session_state.multi_docs)} documents loaded", expanded=False):
            st.markdown(
                "<span style='font-size:0.6rem;font-weight:700;color:rgba(99,102,241,0.55);"
                "text-transform:uppercase;letter-spacing:2px;'>Loaded Documents</span>",
                unsafe_allow_html=True)
            badge_parts = []
            for d in st.session_state.multi_docs:
                is_act = d["name"] == st.session_state.doc_name
                bg = "rgba(99,102,241,0.18)" if is_act else "rgba(30,41,59,0.6)"
                border = "rgba(99,102,241,0.4)" if is_act else "rgba(51,65,85,0.4)"
                color = "#a5b4fc" if is_act else "#64748b"
                badge_parts.append(
                    f"<span style='background:{bg};border:1px solid {border};"
                    f"border-radius:7px;padding:5px 12px;font-size:0.73rem;color:{color};"
                    f"margin:3px 3px;display:inline-block;font-weight:600;'>"
                    f"{'[Active] ' if is_act else ''}{d['name']}</span>"
                )
            st.markdown(
                f"<div style='margin:8px 0 14px;line-height:2;'>{''.join(badge_parts)}</div>",
                unsafe_allow_html=True)
            st.markdown(
                "<span style='font-size:0.6rem;font-weight:700;color:rgba(99,102,241,0.55);"
                "text-transform:uppercase;letter-spacing:2px;margin-bottom:6px;display:block;'>"
                "Ask across all documents:</span>",
                unsafe_allow_html=True)
            mq = st.text_input(
                "Multi-document query",
                key="mq_input",
                placeholder="e.g. What are common themes across all documents?",
                label_visibility="collapsed"
            )
            mc1, mc2 = st.columns(2)
            with mc1:
                if st.button("Query All Documents", width="stretch", key="btn_qall"):
                    if mq and mq.strip():
                        with st.spinner(f"Querying {len(st.session_state.multi_docs)} documents..."):
                            try:
                                ans = lazy("agents.multi_document_agent","query_multiple_documents")(mq, st.session_state.multi_docs)
                            except Exception:
                                ans = query_multi_direct(mq, st.session_state.multi_docs)
                        if not ans or not ans.strip():
                            ans = query_multi_direct(mq, st.session_state.multi_docs)
                        st.session_state.chat.append({"role":"human","content":f"[Multi-Doc Query] {mq}"})
                        st.session_state.chat.append({"role":"assistant","content":ans})
                        st.session_state.pending = {"evidence":[],"intent":{},"question":mq,"answer":ans,"ks":uid()}
                        st.rerun()
                    else:
                        st.warning("Please type a question above before clicking Query All.")
            with mc2:
                if st.button("Compare All Documents", width="stretch", key="btn_cmpall"):
                    with st.spinner(f"Comparing {len(st.session_state.multi_docs)} documents..."):
                        try:
                            comp = lazy("agents.multi_document_agent","compare_documents")(st.session_state.multi_docs)
                        except Exception:
                            comp = compare_docs_direct(st.session_state.multi_docs)
                    if not comp or not comp.strip():
                        comp = compare_docs_direct(st.session_state.multi_docs)
                    st.session_state.chat.append({"role":"human","content":f"Compare all {len(st.session_state.multi_docs)} documents"})
                    st.session_state.chat.append({"role":"assistant","content":comp})
                    st.session_state.pending = {"evidence":[],"intent":{},"question":"compare","answer":comp,"ks":uid()}
                    st.rerun()

    if False and len(st.session_state.multi_docs) > 1:
        with st.expander(f"Multi-Document Panel - {len(st.session_state.multi_docs)} loaded", expanded=False):
            badge_html = "".join([f"<span style='background:rgba({'99,102,241' if d['name']==st.session_state.doc_name else '30,41,59'},0.4);border:1px solid rgba({'99,102,241' if d['name']==st.session_state.doc_name else '51,65,85'},0.3);border-radius:7px;padding:4px 10px;font-size:0.72rem;color:{'#a5b4fc' if d['name']==st.session_state.doc_name else '#64748b'};margin:3px;display:inline-block;'>{'✅' if d['name']==st.session_state.doc_name else '📄'} {d['name']}</span>" for d in st.session_state.multi_docs])
            st.markdown(f"<div style='margin-bottom:12px;'>{badge_html}</div>", unsafe_allow_html=True)
            mq = st.text_input("Ask across all documents:", key="mq_input")
            mc1, mc2 = st.columns(2)
            with mc1:
                if st.button("Query All", width="stretch", key="btn_qall"):
                    if mq:
                        with st.spinner("Querying all..."):
                            try: ans = lazy("agents.multi_document_agent","query_multiple_documents")(mq, st.session_state.multi_docs)
                            except Exception: ans = query_multi_direct(mq, st.session_state.multi_docs)
                        st.session_state.chat.append({"role":"human","content":f"[All Docs] {mq}"})
                        st.session_state.chat.append({"role":"assistant","content":ans})
                        st.session_state.pending = {"evidence":[],"intent":{},"question":mq,"answer":ans,"ks":uid()}; st.rerun()
            with mc2:
                if st.button("Compare All", width="stretch", key="btn_cmpall"):
                    with st.spinner("Comparing..."):
                        try: comp = lazy("agents.multi_document_agent","compare_documents")(st.session_state.multi_docs)
                        except Exception: comp = compare_docs_direct(st.session_state.multi_docs)
                    st.session_state.chat.append({"role":"human","content":"Compare all documents"})
                    st.session_state.chat.append({"role":"assistant","content":comp})
                    st.session_state.pending = {"evidence":[],"intent":{},"question":"compare","answer":comp,"ks":uid()}; st.rerun()

    for idx, msg in enumerate(st.session_state.chat):
        if msg["role"] == "human":
            with st.chat_message("user"): st.markdown(msg["content"])
        else:
            content = msg["content"]
            if content.startswith("__QUIZ__"):
                ks = content.replace("__QUIZ__", "")
                with st.chat_message("assistant"):
                    if ks in st.session_state.quiz_store: render_quiz(st.session_state.quiz_store[ks], ks)
                    else: st.markdown("Quiz expired. Ask again to regenerate.")
            else:
                with st.chat_message("assistant"):
                    st.markdown(content)
                    if idx == len(st.session_state.chat) - 1:
                        render_tts(content, ks=str(idx))

    render_pending(st.session_state.get("pending", {}))

    col_in, col_stop, col_v = st.columns([5, 1, 1])
    with col_in:
        question = st.chat_input("Ask anything about your document...")
        if question:
            st.session_state.pending = {}
            answer_q(question, deps)
            st.rerun()
    with col_stop:
        if st.session_state.get("is_processing"):
            st.markdown("<div class='stop-btn'>", unsafe_allow_html=True)
            if st.button("Stop", width="stretch", key="global_stop"):
                st.session_state.stop_requested = True
                st.session_state.is_processing = False
                st.info("Processing stopped.")
                st.rerun()
            st.markdown("</div>", unsafe_allow_html=True)
    with col_v:
        try: is_https = st.context.headers.get("x-forwarded-proto") == "https"
        except Exception: is_https = False
        if is_https:
            try:
                from streamlit_mic_recorder import mic_recorder
                audio = mic_recorder(start_prompt="🎙", stop_prompt="⏹", key="voice")
                if audio and audio.get("bytes"):
                    fn = lazy("agents.voice_agent", "transcribe_audio_file")
                    vr = fn(audio["bytes"])
                    if vr.get("success"): question = vr.get("text"); st.success(vr.get("text"))
            except Exception: st.button("🎙", help="Voice unavailable")
        else:
            if st.button("🎙", help="Voice requires HTTPS"): st.info("Voice available on Streamlit Cloud.")

    if question:
        st.session_state.pending = {}   # clear old pending NOW
        st.session_state["_regen_done"] = False
        answer_q(question, deps)
        st.rerun()

if __name__ == "__main__":
    main()

